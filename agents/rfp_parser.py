# agents/rfp_parser.py
# STEP 1: RFP 분석 에이전트
# 역할: HWP/HWPX/PDF/TXT 파일에서 텍스트 추출 후 Claude API로 분석
# 출력: 기관정보, 평가항목, 핵심키워드, 과업목록, 금지사항, 톤앤매너 힌트

import json
import re
import struct
import zipfile
import zlib
import xml.etree.ElementTree as ET
from pathlib import Path

from core import claude_client
from core.dna import ConceptDNA, update_dna
from database.db import save_rfp_analysis

# RFP 텍스트 Claude 전달 시 최대 글자 수 (토큰 초과 방지)
_MAX_RFP_CHARS = 12000


# ─────────────────────────────────────────────
# 진입점
# ─────────────────────────────────────────────

def run(dna: ConceptDNA, file_path: str = None) -> dict:
    """RFP 분석 실행.

    Args:
        dna: 사용자 입력이 담긴 초기 ConceptDNA
        file_path: RFP 파일 경로 (HWP/HWPX/PDF/TXT). None이면 dna.rfp_text 사용.

    Returns:
        {
            "basic_info":       {"client_name", "project_name", "budget", "deadline"},
            "agency_type":      str,   # 중앙부처/지자체/의회/공공기관/기타
            "core_tasks":       list,  # 핵심 과업 목록
            "evaluation_items": list,  # [{"item": "...", "score": "..."}]
            "top_keywords":     list,  # 발주처 강조 키워드 TOP 10
            "forbidden_notes":  list,  # 금지/주의 사항
            "agency_tone_hint": str,   # 기관 특성 + 톤앤매너 힌트
        }
    """
    # 1. 텍스트 추출
    if file_path:
        print(f"  파일 텍스트 추출 중: {Path(file_path).name}")
        rfp_text = extract_text(file_path)
    elif dna.rfp_text:
        rfp_text = dna.rfp_text
    else:
        # 파일/텍스트 없으면 사용자 입력값만으로 분석
        rfp_text = ""

    # 2. Claude API로 분석
    print("  Claude API로 RFP 분석 중...")
    result = _analyze_with_claude(rfp_text, dna)

    # 3. DNA 업데이트
    basic = result.get("basic_info", {})
    eval_items = result.get("evaluation_items", [])

    # 평가 배점표를 프롬프트 주입용 포맷 문자열로 변환
    evaluation_criteria = _format_evaluation_criteria(eval_items)

    # 배점 상위 3개 항목명 추출 (top_criteria)
    def _parse_score_int(s: str) -> int:
        m = re.search(r"\d+", str(s or ""))
        return int(m.group()) if m else 0

    sorted_items = sorted(
        [it for it in eval_items if isinstance(it, dict)],
        key=lambda x: _parse_score_int(x.get("score", "")),
        reverse=True
    )
    top_criteria = [it.get("item", "") for it in sorted_items if it.get("item")][:3]

    # 정량 평가 항목 별도 추출 (실적·인력·등급 등 수치 검증 항목)
    quantitative_requirements = [
        it for it in eval_items
        if isinstance(it, dict) and it.get("category", "") == "정량적"
    ]

    # 배점표 전략 분석 추출
    evaluation_strategy = result.get("evaluation_strategy", {})
    if not isinstance(evaluation_strategy, dict):
        evaluation_strategy = {}

    update_dna(dna, {
        "client_name":        basic.get("client_name") or dna.client_name,
        "project_name":       basic.get("project_name") or dna.project_name,
        "budget":             basic.get("budget") or dna.budget,
        "deadline":           basic.get("deadline") or dna.deadline,
        "agency_type":        result.get("agency_type", ""),
        "core_tasks":         result.get("core_tasks", []),
        "evaluation_items":   eval_items,
        "evaluation_criteria": evaluation_criteria,
        "top_criteria":       top_criteria,
        "quantitative_requirements": quantitative_requirements,
        "evaluation_strategy": evaluation_strategy,
        "evaluation_keywords": result.get("top_keywords", []),
        "rfp_requirements":   result.get("core_tasks", []),
        "forbidden_notes":    result.get("forbidden_notes", []),
        "agency_characteristics": result.get("agency_tone_hint", ""),
        "rfp_text":           rfp_text[:2000] if rfp_text else "",
    })

    # 4. DB 저장
    try:
        save_rfp_analysis(dna.client_name, dna.project_name, result,
                          case_id=getattr(dna, "case_id", 0) or 0)
        print("  분석 결과 DB 저장 완료")
    except Exception as e:
        print(f"  [경고] DB 저장 실패 (계속 진행): {e}")

    return result


# ─────────────────────────────────────────────
# 파일 형식별 텍스트 추출
# ─────────────────────────────────────────────
# 폼 자동채우기용 빠른 추출
# ─────────────────────────────────────────────

def rfp_quick_extract(rfp_text: str) -> dict:
    """RFP 텍스트에서 폼 자동채우기 + 전략 분석용 핵심 필드를 추출.

    Returns:
        {
            "client_name": str,
            "project_name": str,
            "budget": str,
            "deadline": str,
            "video_type": str,
            "quantity": int,
            "duration": str,
            "core_keywords": list,   # 핵심 키워드
            "core_tasks": list,      # 핵심 과업
            "prohibited": list,      # 금지 사항
            "special_notes": list,   # 특이사항
            "evaluation_criteria": list,
            "evaluation_strategy": dict,
        }
    """
    text_section = rfp_text[:8000] if len(rfp_text) > 8000 else rfp_text

    prompt = f"""아래 RFP 문서에서 모든 핵심 정보를 추출하라.
없는 항목은 빈 문자열 또는 빈 배열로 반환.

[RFP 문서]
{text_section}

【추출 규칙 — 반드시 준수】
- 평가배점표/심사기준: 모든 항목 빠짐없이 추출, 배점은 정수로 변환
- 핵심 과업: "~을 제작", "~을 수행" 형태의 구체적 과업 항목
- 금지 사항: "금지", "불가", "제한", "제외", "안 됨" 표현 모두
- 특이사항: 감점 조건, 미제출 불이익, 자격 제한
- 핵심 키워드: 발주처가 반복 강조하는 단어/개념 (10개 이내)

아래 JSON만 출력 (다른 텍스트 금지):
{{
  "client_name": "발주처 기관명",
  "project_name": "사업명",
  "budget": "예산 금액",
  "deadline": "납품기한",
  "video_type": "반드시 다음 중 하나: 홍보영상 / 다큐멘터리 / 교육영상 / 캠페인영상 / 뉴스형영상",
  "quantity": 납품 수량 (정수),
  "duration": "편당 러닝타임",
  "core_keywords": ["핵심 키워드 (최대 10개)"],
  "core_tasks": ["구체적 과업 항목"],
  "prohibited": ["금지 또는 제한 사항"],
  "special_notes": ["감점 조건, 자격 제한 등 특이사항"],
  "evaluation_criteria": [
    {{
      "구분": "정성적 또는 정량적 또는 가격 중 하나",
      "항목명": "평가 항목명",
      "배점": 배점 숫자 (정수),
      "세부기준": "단계별 점수 기준 (없으면 빈 문자열)",
      "전략": "이 항목에서 높은 점수 받는 방법 (2문장)",
      "주의사항": "특이사항 (없으면 빈 문자열)"
    }}
  ],
  "evaluation_strategy": {{
    "총점": 전체 배점 합계 숫자,
    "핵심항목": ["배점 10점 이상 항목명"],
    "정량체크리스트": ["항목명: 준비 서류/기준"],
    "집중공략": "가장 점수 올리기 쉬운 항목과 공략법 (3문장)"
  }}
}}
evaluation_criteria: 없으면 빈 배열 []. 배점 높은 순으로 정렬."""

    _empty = {"client_name": "", "project_name": "", "budget": "",
              "deadline": "", "video_type": "", "quantity": 0, "duration": "",
              "core_keywords": [], "core_tasks": [], "prohibited": [], "special_notes": [],
              "evaluation_criteria": [],
              "evaluation_strategy": {"총점": 0, "핵심항목": [], "정량체크리스트": [], "집중공략": ""}}

    try:
        result = claude_client.call_json(prompt, max_tokens=3000, _validate=False)
    except Exception as e:
        # call_json 3단계 폴백 모두 실패 → 원시 응답에서 {} 블록 직접 추출 시도
        print(f"  [rfp_quick_extract] call_json 실패: {e}")
        raw = getattr(e, "__cause__", None)
        raw_text = str(raw) if raw else ""
        m = re.search(r"\{.*\}", raw_text, re.DOTALL)
        if m:
            try:
                from json_repair import repair_json
                result = json.loads(repair_json(m.group()))
            except Exception as e2:
                print(f"  [rfp_quick_extract] 중괄호 블록 파싱도 실패: {e2}")
                return _empty
        else:
            return _empty

    # 타입 보정
    if "quantity" in result:
        try:
            result["quantity"] = int(result["quantity"])
        except (ValueError, TypeError):
            result["quantity"] = 0
    for list_field in ("core_keywords", "core_tasks", "prohibited", "special_notes", "evaluation_criteria"):
        if list_field not in result or not isinstance(result[list_field], list):
            result[list_field] = []
    if "evaluation_strategy" not in result or not isinstance(result["evaluation_strategy"], dict):
        result["evaluation_strategy"] = {"총점": 0, "핵심항목": [], "정량체크리스트": [], "집중공략": ""}
    return result


# ─────────────────────────────────────────────

def extract_text(file_path: str) -> str:
    """파일 형식을 자동 감지하여 텍스트 추출.

    지원 형식: .hwp, .hwpx, .pdf, .txt (대소문자 무관)
    확장자가 없거나 인식 불가 시 파일 내용(매직바이트)으로 형식 자동 감지.

    Args:
        file_path: 파일 경로

    Returns:
        추출된 텍스트 문자열

    Raises:
        ValueError: 지원하지 않는 형식이거나 추출 실패 시
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    extractors = {
        ".txt":  _extract_txt,
        ".pdf":  _extract_pdf,
        ".hwpx": _extract_hwpx,
        ".hwp":  _extract_hwp,
    }

    suffix = path.suffix.lower()

    # 확장자로 추출기 선택 (대소문자 이미 lower() 처리됨)
    extractor = extractors.get(suffix)

    # 확장자 없거나 인식 불가 → 매직바이트로 형식 감지
    if extractor is None:
        detected = _detect_format_by_magic(path)
        if detected:
            print(f"  [자동감지] {path.name!r}: 확장자={suffix!r} → 감지={detected}")
            extractor = extractors[detected]
        else:
            raise ValueError(
                f"지원하지 않는 파일 형식: {suffix!r}  "
                f"(지원: .hwp .hwpx .pdf .txt, 현재 파일: {path.name!r})"
            )

    text = extractor(path)
    if not text or not text.strip():
        raise ValueError(f"파일에서 텍스트를 추출할 수 없습니다: {path.name}")
    return text


def _detect_format_by_magic(path: Path) -> "str | None":
    """파일 첫 바이트로 형식 감지. 인식 못하면 None 반환."""
    try:
        header = path.read_bytes()[:8]
    except Exception:
        return None

    # PDF: %PDF
    if header[:4] == b"%PDF":
        return ".pdf"

    # ZIP 계열 (HWPX = ZIP 구조)
    if header[:4] == b"PK\x03\x04":
        # HWPX인지 확인: 내부에 Contents/content.hpf 가 있으면 HWPX
        try:
            import zipfile as _zf
            with _zf.ZipFile(path) as zf:
                names = zf.namelist()
            if any(n.startswith("Contents/") or n == "mimetype" for n in names):
                return ".hwpx"
        except Exception:
            pass
        return ".hwpx"  # ZIP이면 일단 hwpx 시도

    # HWP 5.x: OLE Compound Document magic
    if header[:8] == b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1":
        return ".hwp"

    # TXT 시도: UTF-8 또는 EUC-KR로 디코딩 가능하면 TXT 처리
    try:
        path.read_bytes()[:512].decode("utf-8")
        return ".txt"
    except Exception:
        pass
    try:
        path.read_bytes()[:512].decode("euc-kr")
        return ".txt"
    except Exception:
        pass

    return None


def _extract_txt(path: Path) -> str:
    """TXT: 인코딩 자동 감지 후 읽기."""
    for encoding in ("utf-8", "euc-kr", "cp949", "utf-8-sig"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    raise ValueError(f"TXT 파일 인코딩을 감지할 수 없습니다: {path.name}")


def _extract_pdf(path: Path) -> str:
    """PDF: pdfplumber로 페이지별 텍스트 추출."""
    try:
        import pdfplumber
    except ImportError:
        raise ImportError("PDF 파싱을 위해 pdfplumber를 설치하세요: pip install pdfplumber")

    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)

    if not pages:
        raise ValueError(f"PDF에서 텍스트를 추출할 수 없습니다 (스캔본일 수 있음): {path.name}")
    return "\n".join(pages)


def _extract_hwpx(path: Path) -> str:
    """HWPX: ZIP 압축 해제 후 section XML 파싱.

    HWPX 구조:
        Contents/section0.xml (또는 section/section*.xml) 에 본문 저장
        텍스트 요소: <hp:t>, <t>, {namespace}t
    """
    if not zipfile.is_zipfile(str(path)):
        raise ValueError(f"유효한 HWPX(ZIP) 파일이 아닙니다: {path.name}")

    text_parts = []
    with zipfile.ZipFile(str(path), "r") as z:
        all_names = z.namelist()

        # section XML 파일 수집 — 다양한 경로 패턴 대응
        section_files = sorted(
            name for name in all_names
            if name.endswith(".xml") and (
                "section" in name.lower() or
                name.startswith("Contents/") or
                name.startswith("content/")
            )
        )

        # section 파일이 없으면 모든 XML 파일 시도
        if not section_files:
            section_files = [n for n in all_names if n.endswith(".xml")]

        if not section_files:
            raise ValueError(f"HWPX 파일에서 XML을 찾을 수 없습니다: {path.name}")

        print(f"  [HWPX] {len(section_files)}개 XML ({section_files[:3]}...)")

        for section_file in section_files:
            try:
                with z.open(section_file) as f:
                    raw_xml = f.read()
            except Exception as e:
                print(f"  [HWPX] {section_file} 열기 실패: {e}")
                continue

            try:
                root = ET.fromstring(raw_xml)
            except ET.ParseError as e:
                print(f"  [HWPX] {section_file} XML 파싱 실패: {e}")
                continue

            for elem in root.iter():
                tag = elem.tag
                # 네임스페이스 제거한 로컬명이 't' 인 요소 (hp:t, {ns}t 등)
                local = tag.split("}")[-1] if "}" in tag else tag
                if local == "t":
                    if elem.text and elem.text.strip():
                        text_parts.append(elem.text.strip())

    if not text_parts:
        raise ValueError(f"HWPX 파일에서 텍스트를 추출할 수 없습니다: {path.name}")
    return "\n".join(text_parts)


def _extract_hwp(path: Path) -> str:
    """HWP 바이너리: olefile로 BodyText 스트림 추출 후 레코드 파싱.

    HWP5 바이너리 구조:
        BodyText/Section0, Section1, ... 스트림에 본문 저장
        각 스트림은 zlib 압축 (wbits=-15)
        레코드 타입 67 (HWPTAG_PARA_TEXT) 에 UTF-16-LE 텍스트 저장
    """
    try:
        import olefile
    except ImportError:
        raise ImportError("HWP 파싱을 위해 olefile을 설치하세요: pip install olefile")

    if not olefile.isOleFile(str(path)):
        raise ValueError(f"유효한 HWP(OLE) 파일이 아닙니다: {path.name}")

    text_parts = []
    with olefile.OleFileIO(str(path)) as ole:
        if not ole.exists("BodyText"):
            raise ValueError(f"HWP 파일에 BodyText 스트림이 없습니다: {path.name}")

        section_idx = 0
        while ole.exists(f"BodyText/Section{section_idx}"):
            data = ole.openstream(f"BodyText/Section{section_idx}").read()
            # zlib 압축 해제 (raw deflate, wbits=-15)
            try:
                decompressed = zlib.decompress(data, -15)
            except zlib.error:
                decompressed = data  # 비압축 섹션 허용

            section_text = _parse_hwp_records(decompressed)
            if section_text:
                text_parts.append(section_text)
            section_idx += 1

    if not text_parts:
        raise ValueError(f"HWP 파일에서 텍스트를 추출하지 못했습니다: {path.name}")
    return "\n".join(text_parts)


def _parse_hwp_records(data: bytes) -> str:
    """HWP 바이너리 레코드 스트림에서 텍스트 파싱.

    HWP 레코드 헤더 구조 (4바이트 little-endian):
        bits  0-9  : record type (tag ID)
        bits 10-19 : level
        bits 20-31 : size (0xFFF이면 다음 4바이트가 실제 크기)

    HWPTAG_PARA_TEXT = 67 : UTF-16-LE 텍스트 데이터
    """
    text_parts = []
    i = 0
    data_len = len(data)

    while i + 4 <= data_len:
        header = struct.unpack_from("<I", data, i)[0]
        rec_type = header & 0x3FF
        rec_size = (header >> 20) & 0xFFF
        i += 4

        # 확장 크기 (0xFFF = 4095 이면 다음 4바이트에 실제 크기)
        if rec_size == 0xFFF:
            if i + 4 > data_len:
                break
            rec_size = struct.unpack_from("<I", data, i)[0]
            i += 4

        if i + rec_size > data_len:
            break

        if rec_type == 67:  # HWPTAG_PARA_TEXT
            chunk = data[i : i + rec_size]
            try:
                text = chunk.decode("utf-16-le", errors="ignore")
                # 제어문자 및 특수 HWP 코드 제거
                text = "".join(
                    ch for ch in text
                    if ch >= " " or ch in ("\n", "\t")
                )
                text = text.strip()
                if text:
                    text_parts.append(text)
            except Exception:
                pass

        i += rec_size

    return "\n".join(text_parts)


# ─────────────────────────────────────────────
# 평가 배점표 포맷 변환
# ─────────────────────────────────────────────

def _format_evaluation_criteria(eval_items: list) -> str:
    """evaluation_items 리스트를 프롬프트 주입용 포맷 문자열로 변환.

    Args:
        eval_items: [{"item", "score", "category", "criteria", "detail_criteria", "warning", "required"}, ...]

    Returns:
        예)
        • 사업이해도 — 20점 [정성적]: 사업목적과 추진배경 이해 수준
        • 유사용역실적 — 5점 [정량적]
          └ 기준: 5건이상=5점, 3건이상=3점
          └ ⚠️ 미제출시 최저점
    """
    if not eval_items:
        return ""
    lines = []
    for it in eval_items:
        if not isinstance(it, dict):
            continue
        name   = (it.get("item") or "").strip()
        score  = (it.get("score") or "").strip()
        cat      = (it.get("category") or "").strip()
        req      = (it.get("required") or "").strip()
        crit     = (it.get("criteria") or "").strip()
        detail   = (it.get("detail_criteria") or "").strip()
        hint     = (it.get("strategic_hint") or "").strip()
        warn     = (it.get("warning") or "").strip()
        if not name:
            continue
        parts = [f"• {name}"]
        if score:
            parts.append(f"— {score}")
        if cat:
            parts.append(f"[{cat}]")
        if req:
            parts.append(f"({req})")
        if crit:
            parts.append(f": {crit}")
        line = " ".join(parts)
        if detail:
            line += f"\n  └ 기준: {detail}"
        if hint:
            line += f"\n  └ 전략: {hint}"
        if warn:
            line += f"\n  └ ⚠️ {warn}"
        lines.append(line)
    return "\n".join(lines)


# ─────────────────────────────────────────────
# Claude API 분석
# ─────────────────────────────────────────────

def _analyze_with_claude(rfp_text: str, dna: ConceptDNA) -> dict:
    """추출된 텍스트 + 사용자 입력으로 RFP 분석."""
    prompt = _build_prompt(rfp_text, dna)
    return claude_client.call_json(prompt, max_tokens=4096)


def _build_prompt(rfp_text: str, dna: ConceptDNA) -> str:
    """RFP 분석용 프롬프트 생성."""
    # 텍스트가 너무 길면 앞부분 우선 사용 (평가항목은 보통 앞에 위치)
    if len(rfp_text) > _MAX_RFP_CHARS:
        rfp_section = rfp_text[:_MAX_RFP_CHARS] + "\n...(이하 생략)"
    elif rfp_text:
        rfp_section = rfp_text
    else:
        rfp_section = "(RFP 문서 없음 — 사용자 입력값만으로 분석)"

    user_inputs = []
    if dna.client_name:
        user_inputs.append(f"- 발주처: {dna.client_name}")
    if dna.project_name:
        user_inputs.append(f"- 사업명: {dna.project_name}")
    if dna.video_type:
        user_inputs.append(f"- 영상 종류: {dna.video_type}")
    if dna.budget:
        user_inputs.append(f"- 예산: {dna.budget}")
    if dna.deadline:
        user_inputs.append(f"- 납품기한: {dna.deadline}")
    user_input_block = "\n".join(user_inputs) if user_inputs else "(없음)"

    return f"""당신은 대한민국 정부 입찰 전문가입니다.
아래 RFP(제안요청서) 문서와 사용자 입력값을 분석해서 JSON으로만 결과를 반환하세요.

[사용자 입력값]
{user_input_block}

[RFP 문서]
{rfp_section}

[분석 지침]
7개 항목을 분석하세요. 문서에서 명확히 확인되지 않는 항목은 사용자 입력값과 문맥으로 추론하세요.

1. basic_info — 기관명, 사업명, 예산(금액+단위), 납품기한 추출
2. agency_type — 반드시 다음 중 하나: 중앙부처 / 지자체 / 의회 / 공공기관 / 기타
3. core_tasks — 핵심 과업 목록 (구체적 수행 항목, 배열)
4. evaluation_items — 평가 배점표 (배열). 평가항목/배점표가 있으면 반드시 모두 추출하라.
   각 항목:
   {{
     "item": "항목명",
     "score": "배점(숫자+점, 예: 20점)",
     "category": "정성적 / 정량적 / 가격 중 하나",
     "criteria": "평가 기준 설명 (없으면 빈 문자열)",
     "detail_criteria": "단계별 점수 기준 (예: 5건이상=5점, 3건이상=3점. 없으면 빈 문자열)",
     "strategic_hint": "이 항목에서 높은 점수를 받으려면 제안서에 구체적으로 무엇을 어떻게 써야 하는지 (2~3문장)",
     "warning": "미제출시 최저점 등 특이사항 (없으면 빈 문자열)",
     "required": "필수/선택 (없으면 빈 문자열)",
     "importance": "배점 기준 중요도 — high(배점 상위 30%) / medium / low"
   }}
   예) {{"item": "유사용역수행실적", "score": "5점", "category": "정량적", "criteria": "유사 용역 수행 실적 평가", "detail_criteria": "5건이상=5점, 3건이상=3점, 1건이상=1점", "strategic_hint": "유사 공공기관 영상 제작 실적 5건 이상을 실적증명서와 계약서로 제출. 3천만원 이상 계약 건을 우선 제시.", "warning": "미제출시 최저점", "required": "필수", "importance": "medium"}}
   배점이 없으면 score: "" 로 입력
5. evaluation_strategy — 배점표 기반 전략 분석 (단일 객체)
   {{
     "총점": 전체 배점 합계 (정수),
     "핵심항목": ["배점 10점 이상 항목명 리스트"],
     "정량항목_체크리스트": ["각 정량 항목별 '항목명: 준비 서류/기준' 형식으로"],
     "집중공략": "가장 점수 올리기 쉬운 항목과 구체적 공략법 (3~4문장)"
   }}
5. top_keywords — 발주처가 문서에서 강조하는 핵심 키워드 TOP 10 (단어 또는 짧은 구, 배열)
6. forbidden_notes — 금지/주의 사항 목록 (배열, 없으면 빈 배열)
7. agency_tone_hint — 기관 특성 요약 및 톤앤매너 힌트 (2~3문장)

반드시 아래 JSON 형식으로만 출력하세요. 다른 설명은 하지 마세요:
{{
  "basic_info": {{
    "client_name": "...",
    "project_name": "...",
    "budget": "...",
    "deadline": "..."
  }},
  "agency_type": "...",
  "core_tasks": ["...", "..."],
  "evaluation_items": [
    {{
      "item": "...", "score": "...", "category": "정성적/정량적/가격",
      "criteria": "...", "detail_criteria": "...", "strategic_hint": "...",
      "warning": "...", "required": "...", "importance": "high/medium/low"
    }}
  ],
  "evaluation_strategy": {{
    "총점": 0,
    "핵심항목": ["..."],
    "정량항목_체크리스트": ["항목명: 준비 사항"],
    "집중공략": "..."
  }},
  "top_keywords": ["...", "..."],
  "forbidden_notes": ["...", "..."],
  "agency_tone_hint": "..."
}}"""


# ─────────────────────────────────────────────
# 참고 제안서 분석
# ─────────────────────────────────────────────

def parse_reference_proposal(file_path: str) -> str:
    """참고 제안서 파일을 분석해 구조·흐름 요약 반환.

    Args:
        file_path: 참고 제안서 파일 경로 (PDF/PPTX/DOCX/HWP/HWPX)

    Returns:
        구조 분석 요약 문자열 (DNA.reference_structure에 저장됨)
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    # 확장자 없거나 인식 불가 시 매직바이트로 감지
    if not suffix or suffix not in (".hwp", ".hwpx", ".pdf", ".txt", ".docx", ".pptx"):
        detected = _detect_format_by_magic(path)
        if detected:
            print(f"  [참고 제안서 자동감지] {path.name!r}: {detected}")
            suffix = detected
        else:
            print(f"  [참고 제안서] 형식 감지 실패: {path.name!r}")
            return ""

    # 텍스트 추출
    text = ""
    try:
        if suffix in (".hwp", ".hwpx", ".pdf", ".txt"):
            text = extract_text(file_path)
        elif suffix in (".docx",):
            text = _extract_docx(path)
        elif suffix in (".pptx",):
            text = _extract_pptx_text(path)
        else:
            print(f"  [참고 제안서] 지원하지 않는 형식: {suffix!r}")
            return ""
    except Exception as e:
        print(f"  [참고 제안서] 텍스트 추출 실패: {e}")
        return ""

    if not text or not text.strip():
        print("  [참고 제안서] 텍스트 추출 결과 없음")
        return ""

    # 텍스트 길이 제한
    text_excerpt = text[:10000]
    print(f"  [참고 제안서] 분석 중... ({len(text_excerpt)}자)")

    prompt = f"""당신은 제안서 구조 분석 전문가입니다.
아래는 과거 낙찰 제안서(또는 참고 제안서)의 텍스트입니다.
이 제안서의 구조와 설득 패턴을 분석해 아래 6가지 항목을 JSON으로 작성하세요.

[제안서 텍스트]
{text_excerpt}

반드시 아래 JSON 형식으로만 출력하세요:
{{
  "toc_structure": "목차 구성 (섹션 순서를 번호 목록으로)",
  "persuasion_flow": "설득 흐름 설명 (어떤 논리 순서로 발주처를 설득하는지 3~4문장)",
  "volume_distribution": "분량 배분 설명 (어느 섹션에 얼마나 할애했는지)",
  "tone_and_style": "문체/톤앤매너 설명 (격식체/구어체, 수치 활용도, 시각화 방식 등)",
  "differentiation_method": "차별화 포인트 제시 방식 (어떻게 경쟁사 대비 우위를 표현했는지)",
  "evidence_method": "실적/증거 제시 방식 (포트폴리오, 수상실적, 레퍼런스 활용 방법)"
}}"""

    try:
        result = claude_client.call_json(prompt, max_tokens=3000)
        # 6개 항목을 읽기 쉬운 텍스트로 변환
        lines = [
            f"[목차 구성] {result.get('toc_structure', '')}",
            f"[설득 흐름] {result.get('persuasion_flow', '')}",
            f"[분량 배분] {result.get('volume_distribution', '')}",
            f"[문체/톤앤매너] {result.get('tone_and_style', '')}",
            f"[차별화 방식] {result.get('differentiation_method', '')}",
            f"[실적 제시] {result.get('evidence_method', '')}",
        ]
        summary = "\n".join(lines)
        print(f"  [참고 제안서] 분석 완료")
        return summary
    except Exception as e:
        print(f"  [참고 제안서] Claude 분석 실패: {e}")
        return ""


def _extract_docx(path: Path) -> str:
    """DOCX 텍스트 추출."""
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        raise ImportError("DOCX 파싱을 위해 python-docx를 설치하세요: pip install python-docx")


def _extract_pptx_text(path: Path) -> str:
    """PPTX 슬라이드 텍스트 추출."""
    try:
        from pptx import Presentation as _Prs
        prs = _Prs(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = " ".join(run.text for run in para.runs if run.text.strip())
                        if text.strip():
                            parts.append(text)
        return "\n".join(parts)
    except ImportError:
        raise ImportError("PPTX 파싱을 위해 python-pptx를 설치하세요")
