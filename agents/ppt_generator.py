# agents/ppt_generator.py
# PPT 생성: Claude로 슬라이드 구성 → python-pptx로 파일 생성
# 디자인: 흰 배경 / 검정 텍스트 / 검정 테두리 사각형 / 컬러 없음

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

from core.claude_client import call_json

# ── 모노크롬 색상
_BLACK = RGBColor(0x00, 0x00, 0x00)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GRAY  = RGBColor(0x77, 0x77, 0x77)
_LGRAY = RGBColor(0xCC, 0xCC, 0xCC)

# ── 슬라이드 크기 (와이드스크린 16:9)
W = Inches(13.33)
H = Inches(7.5)
M = Inches(0.55)   # 기본 여백


# ─────────────────────────────────────────────
# 저수준 그리기 헬퍼
# ─────────────────────────────────────────────

def _border_rect(slide, left, top, width, height, fill_rgb=None, border_rgb=_BLACK, border_pt=1.2):
    """테두리 사각형. fill_rgb=None 이면 흰 배경."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _WHITE
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(border_pt)
    return shape


def _no_border_rect(slide, left, top, width, height, fill_rgb=_WHITE):
    """테두리 없는 사각형 (구분선 등에 사용)."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def _txt(slide, left, top, width, height, text,
         size=16, bold=False, color=_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    """텍스트박스 추가."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color
    return tx


def _pagenum(slide, num, total):
    _txt(slide, W - Inches(1.6), H - Inches(0.43),
         Inches(1.4), Inches(0.35),
         f"{num}  /  {total}", size=11, color=_GRAY, align=PP_ALIGN.RIGHT)


def _hline(slide, left, top, width, thickness=Inches(0.015), color=_BLACK):
    """수평 구분선."""
    _no_border_rect(slide, left, top, width, thickness, fill_rgb=color)


# ─────────────────────────────────────────────
# 슬라이드 유형별 그리기 (7종)
# ─────────────────────────────────────────────

def _draw_cover(slide, sl, num, total, meta):
    """1. 표지: 중앙 대형 제목 + 클라이언트 + 날짜"""
    title   = sl.get("title", "")
    client  = sl.get("client", meta.get("client_name", ""))
    date    = sl.get("date", "")

    # 상단 검정 띠 (얇은 강조선)
    _no_border_rect(slide, 0, 0, W, Inches(0.08), fill_rgb=_BLACK)
    # 하단 검정 띠
    _no_border_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill_rgb=_BLACK)

    # 클라이언트
    if client:
        _txt(slide, M, Inches(2.1), W - M*2, Inches(0.5),
             client, size=15, color=_GRAY, align=PP_ALIGN.CENTER)

    # 메인 타이틀
    _txt(slide, M, Inches(2.75), W - M*2, Inches(1.8),
         title, size=36, bold=True, color=_BLACK, align=PP_ALIGN.CENTER)

    # 중앙 구분선
    _hline(slide, Inches(4.0), Inches(4.75), Inches(5.33))

    # 날짜
    if date:
        _txt(slide, M, Inches(5.0), W - M*2, Inches(0.5),
             date, size=13, color=_GRAY, align=PP_ALIGN.CENTER)

    _pagenum(slide, num, total)


def _draw_toc(slide, sl, num, total):
    """2. 목차: 번호 붙은 섹션 목록"""
    title   = sl.get("title", "목차")
    items   = sl.get("items", [])

    # 제목 영역
    _txt(slide, M, Inches(0.45), W - M*2, Inches(0.7),
         title, size=22, bold=True, color=_BLACK)
    _hline(slide, M, Inches(1.25), W - M*2)

    # 항목 목록
    row_h  = Inches(0.65)
    start_y = Inches(1.5)
    col_w  = (W - M*2) / 2
    for idx, item in enumerate(items[:10]):
        row    = idx // 2
        col    = idx % 2
        x      = M + col * col_w
        y      = start_y + row * row_h
        label  = f"{idx + 1:02d}.  {item}" if isinstance(item, str) else f"{idx + 1:02d}.  {item.get('title', '')}"
        _txt(slide, x, y, col_w - Inches(0.2), row_h,
             label, size=15, color=_BLACK)

    _pagenum(slide, num, total)


def _draw_content(slide, sl, num, total):
    """3. 내용: 제목 바 + 3~5 불릿 포인트"""
    title   = sl.get("title", "")
    bullets = sl.get("bullets", sl.get("content", []))
    if isinstance(bullets, str):
        bullets = [bullets]

    # 제목 바 (검정 배경, 흰 글자)
    TH = Inches(1.0)
    _no_border_rect(slide, 0, 0, W, TH, fill_rgb=_BLACK)
    _txt(slide, M, Inches(0.18), W - M*2, Inches(0.7),
         title, size=20, bold=True, color=_WHITE)

    # 불릿 목록
    body_top = TH + Inches(0.25)
    body_h   = H - body_top - Inches(0.55)
    if bullets:
        bullet_txt = "\n".join(f"•   {line}" for line in bullets if line)
        _txt(slide, M, body_top, W - M*2, body_h,
             bullet_txt, size=16, color=_BLACK)

    _pagenum(slide, num, total)


def _draw_process(slide, sl, num, total):
    """4. 프로세스: [Step1] → [Step2] → [Step3] → [Step4]"""
    title  = sl.get("title", "")
    steps  = sl.get("steps", [])

    # 제목
    _txt(slide, M, Inches(0.45), W - M*2, Inches(0.6),
         title, size=20, bold=True, color=_BLACK)
    _hline(slide, M, Inches(1.15), W - M*2)

    n = len(steps)
    if n == 0:
        _pagenum(slide, num, total)
        return

    # 박스 배치 계산
    box_area_w = W - M*2
    arrow_w    = Inches(0.45)
    box_w      = (box_area_w - arrow_w * (n - 1)) / n
    box_h      = Inches(1.8)
    box_y      = Inches(2.85)

    for idx, step in enumerate(steps):
        label = step if isinstance(step, str) else step.get("label", f"Step {idx+1}")
        desc  = "" if isinstance(step, str) else step.get("desc", "")
        bx    = M + idx * (box_w + arrow_w)

        # 박스
        _border_rect(slide, bx, box_y, box_w, box_h)

        # 스텝 번호
        _txt(slide, bx, box_y + Inches(0.15), box_w, Inches(0.45),
             f"STEP {idx+1}", size=11, bold=True, color=_GRAY, align=PP_ALIGN.CENTER)

        # 스텝 라벨
        _txt(slide, bx + Inches(0.1), box_y + Inches(0.55), box_w - Inches(0.2), Inches(0.65),
             label, size=14, bold=True, color=_BLACK, align=PP_ALIGN.CENTER)

        # 스텝 설명
        if desc:
            _txt(slide, bx + Inches(0.1), box_y + Inches(1.2), box_w - Inches(0.2), Inches(0.55),
                 desc, size=11, color=_GRAY, align=PP_ALIGN.CENTER)

        # 화살표 (마지막 박스 이후엔 없음)
        if idx < n - 1:
            ax = bx + box_w + Inches(0.05)
            ay = box_y + box_h / 2 - Inches(0.02)
            # 화살표 선
            _no_border_rect(slide, ax, ay, arrow_w - Inches(0.12), Inches(0.04), fill_rgb=_BLACK)
            # 화살표 머리 (▶ 텍스트로 대체)
            _txt(slide, ax + arrow_w - Inches(0.3), ay - Inches(0.18),
                 Inches(0.3), Inches(0.4), "▶", size=12, color=_BLACK, align=PP_ALIGN.CENTER)

    _pagenum(slide, num, total)


def _draw_compare(slide, sl, num, total):
    """5. 비교: [좌측 박스: 문제] | [우측 박스: 해결]"""
    title       = sl.get("title", "")
    left_title  = sl.get("left_title", "현재 문제")
    left_items  = sl.get("left_items", [])
    right_title = sl.get("right_title", "해결 방향")
    right_items = sl.get("right_items", [])

    # 슬라이드 제목
    _txt(slide, M, Inches(0.45), W - M*2, Inches(0.6),
         title, size=20, bold=True, color=_BLACK)
    _hline(slide, M, Inches(1.15), W - M*2)

    # 박스 크기
    half_w  = (W - M*2 - Inches(0.3)) / 2
    box_top = Inches(1.4)
    box_h   = H - box_top - Inches(0.6)
    lx      = M
    rx      = M + half_w + Inches(0.3)

    # 좌측 박스
    _border_rect(slide, lx, box_top, half_w, box_h)
    _txt(slide, lx + Inches(0.2), box_top + Inches(0.2), half_w - Inches(0.4), Inches(0.5),
         left_title, size=15, bold=True, color=_BLACK, align=PP_ALIGN.CENTER)
    _hline(slide, lx + Inches(0.15), box_top + Inches(0.8), half_w - Inches(0.3), color=_LGRAY)
    if left_items:
        body = "\n".join(f"•  {it}" for it in left_items if it)
        _txt(slide, lx + Inches(0.2), box_top + Inches(0.95), half_w - Inches(0.4), box_h - Inches(1.1),
             body, size=14, color=_BLACK)

    # 중앙 구분선
    cx = M + half_w + Inches(0.12)
    cy = box_top + Inches(0.3)
    _no_border_rect(slide, cx, cy, Inches(0.06), box_h - Inches(0.6), fill_rgb=_BLACK)

    # 우측 박스
    _border_rect(slide, rx, box_top, half_w, box_h)
    _txt(slide, rx + Inches(0.2), box_top + Inches(0.2), half_w - Inches(0.4), Inches(0.5),
         right_title, size=15, bold=True, color=_BLACK, align=PP_ALIGN.CENTER)
    _hline(slide, rx + Inches(0.15), box_top + Inches(0.8), half_w - Inches(0.3), color=_LGRAY)
    if right_items:
        body = "\n".join(f"•  {it}" for it in right_items if it)
        _txt(slide, rx + Inches(0.2), box_top + Inches(0.95), half_w - Inches(0.4), box_h - Inches(1.1),
             body, size=14, color=_BLACK)

    _pagenum(slide, num, total)


def _draw_number(slide, sl, num, total):
    """6. 숫자 강조: 대형 숫자(48pt) 중앙 + 설명(14pt)"""
    title       = sl.get("title", "")
    number      = str(sl.get("value", "") or sl.get("number", ""))
    label       = sl.get("label", "")
    description = sl.get("description", "")

    # 슬라이드 제목
    _txt(slide, M, Inches(0.45), W - M*2, Inches(0.6),
         title, size=20, bold=True, color=_BLACK)
    _hline(slide, M, Inches(1.15), W - M*2)

    # 중앙 대형 숫자
    _txt(slide, M, Inches(2.0), W - M*2, Inches(1.8),
         number, size=72, bold=True, color=_BLACK, align=PP_ALIGN.CENTER)

    # 숫자 라벨
    if label:
        _txt(slide, M, Inches(3.85), W - M*2, Inches(0.5),
             label, size=16, bold=True, color=_BLACK, align=PP_ALIGN.CENTER)

    # 설명
    if description:
        _hline(slide, Inches(4.0), Inches(4.55), Inches(5.33), color=_LGRAY)
        _txt(slide, M, Inches(4.8), W - M*2, Inches(0.8),
             description, size=14, color=_GRAY, align=PP_ALIGN.CENTER)

    _pagenum(slide, num, total)


def _draw_message(slide, sl, num, total):
    """7. 핵심 메시지: 대형 중앙 문장(28pt) + 하단 노트"""
    title   = sl.get("title", "")
    message = sl.get("message", "")
    note    = sl.get("note", "")

    # 슬라이드 제목
    if title:
        _txt(slide, M, Inches(0.45), W - M*2, Inches(0.6),
             title, size=20, bold=True, color=_BLACK)
        _hline(slide, M, Inches(1.15), W - M*2)

    # 중앙 메시지 박스
    msg_top = Inches(1.8) if title else Inches(2.1)
    _border_rect(slide, M, msg_top, W - M*2, Inches(2.2))
    _txt(slide, M + Inches(0.3), msg_top + Inches(0.35),
         W - M*2 - Inches(0.6), Inches(1.5),
         message, size=28, bold=True, color=_BLACK, align=PP_ALIGN.CENTER)

    # 하단 노트
    if note:
        note_top = msg_top + Inches(2.45)
        _txt(slide, M, note_top, W - M*2, Inches(0.6),
             note, size=13, color=_GRAY, align=PP_ALIGN.CENTER)

    _pagenum(slide, num, total)


# ─────────────────────────────────────────────
# Claude 프롬프트용 케이스 요약
# ─────────────────────────────────────────────

def _build_case_summary(case_detail: dict) -> str:
    case  = case_detail.get("case", {})
    steps = case_detail.get("steps", {})
    dna   = case.get("dna", {})

    lines = [
        f"# {case.get('client_name','')} / {case.get('project_name','')}",
        f"영상 종류: {case.get('video_type','')} | 예산: {case.get('budget','')} | 납품기한: {case.get('deadline','')}",
        "",
    ]

    # ── 컨셉 & 슬로건 (핵심 관통 메시지)
    cr = steps.get("creative", {})
    concept = cr.get("concept", "") or dna.get("concept", "")
    slogan  = cr.get("confirmed_slogan", "") or dna.get("slogan", "")
    if concept or slogan:
        lines.append("## 핵심 컨셉 & 슬로건")
        if concept:
            lines.append(f"컨셉: {concept[:300]}")
            desc = cr.get("concept_description", "") or dna.get("concept_description", "")
            if desc:
                lines.append(f"컨셉 설명: {str(desc)[:300]}")
        if slogan:
            lines.append(f"슬로건: {slogan}")
        tone = cr.get("tone_description", "") or dna.get("tone_and_manner", "")
        if tone:
            lines.append(f"톤앤매너: {str(tone)[:200]}")
        lines.append("")

    # ── 평가 배점표
    rfp = steps.get("rfp_analysis", {})
    eval_criteria = dna.get("evaluation_criteria", "") or rfp.get("evaluation_criteria", "")
    if not eval_criteria:
        eval_items = rfp.get("evaluation_items", []) or dna.get("evaluation_items", [])
        if eval_items:
            eval_criteria = "\n".join(
                f"  - {it.get('item','') if isinstance(it, dict) else str(it)}"
                + (f" ({it['score']}점)" if isinstance(it, dict) and it.get('score') else "")
                for it in eval_items[:12]
            )
    if eval_criteria:
        lines.append("## 평가 배점표 (배점 높은 순 — 슬라이드 분량 배분 기준)")
        lines.append(str(eval_criteria)[:1200])
        lines.append("")

    # ── 핵심 과업 & 키워드
    core_tasks = rfp.get("core_tasks", []) or dna.get("core_tasks", [])
    top_kw     = rfp.get("top_keywords", []) or dna.get("evaluation_keywords", [])
    if core_tasks:
        lines.append("## 핵심 과업")
        lines += [f"  - {str(t)[:150]}" for t in core_tasks[:8]]
        lines.append("")
    if top_kw:
        lines.append("## 평가 핵심 키워드")
        lines.append("  " + ", ".join(str(k)[:40] for k in top_kw[:12]))
        lines.append("")

    # ── 전략
    strat = steps.get("strategy", {})
    if strat:
        lines.append("## 전략")
        for k, lbl in [("core_problem","핵심문제"), ("crisis_statement","위기제시"),
                        ("current_situation","현황진단"), ("solution_direction","해결방향")]:
            if strat.get(k):
                lines.append(f"{lbl}: {str(strat[k])[:280]}")
        effects = strat.get("expected_effects", [])
        if effects:
            lines.append("기대효과: " + " / ".join(str(e)[:100] for e in effects[:5]))

        # 설득 구조 (4단계)
        ps = strat.get("persuasion_structure", [])
        if ps:
            lines.append("설득 구조:")
            for stage in ps[:5]:
                if isinstance(stage, dict):
                    sname = stage.get("stage", "")
                    sbody = stage.get("body", "") or stage.get("description", "")
                    lines.append(f"  [{sname}] {str(sbody)[:180]}")
                elif isinstance(stage, str):
                    lines.append(f"  - {stage[:180]}")

        # 배점 상위 평가항목
        hi = strat.get("high_priority_eval", []) or strat.get("high_priority_eval_items", [])
        if hi:
            lines.append("배점 상위 항목: " + " / ".join(
                (it.get("item","") if isinstance(it, dict) else str(it))[:60]
                for it in hi[:5]
            ))
        lines.append("")

    # ── 리서치 인사이트
    research = steps.get("research", {})
    if research:
        lines.append("## 리서치 인사이트")
        issues = research.get("recent_issues", [])
        if issues:
            lines.append("최근 이슈:")
            for iss in issues[:4]:
                if isinstance(iss, dict):
                    lines.append(f"  - {str(iss.get('title','') or iss.get('issue',''))[:150]}")
                elif isinstance(iss, str):
                    lines.append(f"  - {iss[:150]}")
        sim = research.get("similar_cases", [])
        if sim:
            lines.append("유사 사례:")
            for sc in sim[:3]:
                if isinstance(sc, dict):
                    lines.append(f"  - {str(sc.get('title','') or sc.get('case',''))[:150]}")
        lines.append("")

    # ── 제작 계획
    plan = steps.get("plan", {})
    if plan:
        lines.append("## 제작 계획")
        for ep in plan.get("episodes", [])[:6]:
            if isinstance(ep, dict):
                lines.append(
                    f"  {ep.get('episode_number','')}편: {ep.get('title','')} "
                    f"— {str(ep.get('core_message',''))[:120]}"
                )
        sched = plan.get("production_schedule", [])
        if sched:
            lines.append("제작 일정:")
            for ph in sched[:4]:
                if isinstance(ph, dict):
                    lines.append(f"  [{ph.get('phase','')}] {str(ph.get('tasks',''))[:120]}")
        lines.append("")

    # ── 마케팅
    mkt = steps.get("marketing", {})
    if mkt:
        lines.append("## 마케팅 전략")
        pl = mkt.get("platforms", [])
        if pl:
            lines.append("채널: " + ", ".join(str(p)[:40] for p in pl[:6]))
        for k, lbl in [("target_audience","타겟"), ("kpi","KPI"), ("budget_allocation","예산배분")]:
            if mkt.get(k):
                lines.append(f"{lbl}: {str(mkt[k])[:200]}")
        lines.append("")

    return "\n".join(lines)


# ─────────────────────────────────────────────
# Claude → 슬라이드 JSON 생성
# ─────────────────────────────────────────────

def generate_slides(case_detail: dict, pages: int, progress_cb=None) -> dict:
    if progress_cb:
        progress_cb("Claude AI로 슬라이드 구성 생성 중...", 0, pages)

    summary  = _build_case_summary(case_detail)
    case     = case_detail.get("case", {})
    dna      = case.get("dna", {})
    steps    = case_detail.get("steps", {})
    client   = case.get("client_name", "")

    cr      = steps.get("creative", {})
    concept = cr.get("concept", "") or dna.get("concept", "")
    slogan  = cr.get("confirmed_slogan", "") or dna.get("slogan", "")

    import datetime
    today = datetime.date.today().strftime("%Y년 %m월")

    prompt = f"""당신은 정부 제안 PT 전문 기획자입니다. 아래 제안서 데이터를 분석해 {pages}페이지 PPT 슬라이드 구성을 JSON으로 출력하세요.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[제안서 데이터]
{summary}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

【핵심 관통 메시지 — 표지·목차·마지막 슬라이드에 반드시 반영】
컨셉: {concept or '(데이터에서 추출)'}
슬로건: {slogan or '(데이터에서 추출)'}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
【슬라이드 작성 5원칙】

1. HEAD COPY + SUB COPY + 데이터/출처 구조
   · title(헤드카피): 해당 슬라이드의 핵심 주장을 직관적 한 문장으로 (동사 포함, 15자 이내)
     예시 ✗ "전략 방향"  ✓ "선택과 집중으로 예산 효율 30% 확보"
   · bullets / steps / message(서브카피): 헤드카피를 뒷받침하는 구체적 근거 3~5개
   · note / desc 필드: 수치·출처·데이터 명시 (가능한 경우)

2. 내러티브 연속성
   · 전체 슬라이드가 하나의 스토리를 형성: 문제 제시 → 근거 → 해결책 → 실행 → 기대효과 → 마무리
   · 각 슬라이드 헤드카피를 순서대로 읽으면 논리적 흐름이 형성되어야 함

3. 평가 배점 기반 슬라이드 분량 배분
   · 평가 배점표의 배점이 높은 항목일수록 더 많은 슬라이드 페이지 할당
   · 배점 상위 3개 항목은 각각 최소 2페이지 이상 구성
   · 배점이 낮은 항목은 1페이지로 압축

4. 컨셉·슬로건 관통
   · 표지(cover) title에 슬로건 또는 컨셉 핵심어 반영
   · 목차(toc) 마지막 줄에 슬로건 배치
   · 마지막 슬라이드(message) message 필드에 슬로건 전문 배치

5. INTERZ 차별화 전략 3가지 명시
   · 경쟁사 대비 차별점을 제안서 중반부에 슬라이드 1장으로 구성
   · type: compare 사용 — left: 일반적 접근, right: INTERZ(인터즈)만의 접근
   · INTERZ의 강점(크리에이티브·전략·데이터 기반 실행력)을 근거로 작성

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
【슬라이드 타입 7종 스펙】
  · cover  : title(헤드카피·슬로건 반영), client, date
  · toc    : title, items(섹션 목록, 마지막에 슬로건 1줄)
  · content: title(헤드카피), bullets(서브카피 3~5개)
  · process: title(헤드카피), steps=[{{label, desc}}] (3~5단계)
  · compare: title(헤드카피), left_title, left_items(3~5개), right_title, right_items(3~5개)
  · number : title(헤드카피), number(수치), label(단위/지표명), description(출처·맥락)
  · message: title(선택), message(슬로건·핵심문장 1~2줄), note(보조설명)

【필수 흐름 (이 순서를 지킬 것)】
표지 → 목차 → 현황/문제 → 전략방향 → 크리에이티브컨셉 → INTERZ차별화 → 제작계획 → 마케팅 → 기대효과 → 마무리

반드시 아래 JSON만 출력하세요 (설명 텍스트·마크다운 없이 순수 JSON만):
{{
  "slides": [
    {{
      "number": 1,
      "type": "cover",
      "title": "슬로건 또는 컨셉 반영 제안서 타이틀",
      "client": "{client}",
      "date": "{today}"
    }},
    {{
      "number": 2,
      "type": "toc",
      "title": "목차",
      "items": ["01. 현황 진단 및 문제 정의", "02. 전략 방향", "03. 크리에이티브 컨셉", "04. 인터즈만의 차별화", "05. 제작 계획", "06. 마케팅 전략", "07. 기대 효과", "{slogan or '슬로건'}"]
    }},
    {{
      "number": 3,
      "type": "compare",
      "title": "헤드카피: 현황과 해결 방향 한 문장",
      "left_title": "현재 문제",
      "left_items": ["문제 근거 1 (수치 포함)", "문제 근거 2", "문제 근거 3"],
      "right_title": "해결 방향",
      "right_items": ["해결책 1", "해결책 2", "해결책 3"]
    }},
    {{
      "number": 4,
      "type": "content",
      "title": "헤드카피: 전략 핵심 주장 한 문장",
      "bullets": ["전략 서브포인트 1", "전략 서브포인트 2", "전략 서브포인트 3"]
    }},
    {{
      "number": 5,
      "type": "message",
      "title": "크리에이티브 컨셉",
      "message": "컨셉 핵심 문장 (슬로건 반영)",
      "note": "컨셉 설명 보조 문구"
    }},
    {{
      "number": 6,
      "type": "compare",
      "title": "일반 대행사 vs INTERZ: 무엇이 다른가",
      "left_title": "일반적 접근",
      "left_items": ["차별화 포인트 1 반대 사례", "차별화 포인트 2 반대 사례", "차별화 포인트 3 반대 사례"],
      "right_title": "INTERZ만의 방식",
      "right_items": ["INTERZ 차별점 1 (크리에이티브)", "INTERZ 차별점 2 (전략)", "INTERZ 차별점 3 (데이터 기반 실행)"]
    }},
    {{
      "number": 7,
      "type": "process",
      "title": "헤드카피: 제작 프로세스 핵심 메시지",
      "steps": [
        {{"label": "단계명", "desc": "단계 설명"}},
        {{"label": "단계명", "desc": "단계 설명"}},
        {{"label": "단계명", "desc": "단계 설명"}},
        {{"label": "단계명", "desc": "단계 설명"}}
      ]
    }},
    {{
      "number": 8,
      "type": "number",
      "title": "헤드카피: 기대 효과 수치 핵심 주장",
      "number": "수치",
      "label": "지표명 (단위 포함)",
      "description": "수치 출처 또는 달성 조건"
    }},
    {{
      "number": 9,
      "type": "message",
      "title": "",
      "message": "{slogan or '슬로건 전문'}",
      "note": "제안사: INTERZ(인터즈)"
    }}
  ]
}}

【최종 검증 체크리스트 — 출력 전 반드시 확인】
□ slides 배열이 정확히 {pages}개인가?
□ 1번은 cover, 2번은 toc인가?
□ 마지막은 message 타입이며 슬로건이 포함되어 있는가?
□ INTERZ 차별화 슬라이드(compare 타입)가 1장 이상 있는가?
□ 모든 title(헤드카피)이 "섹션명"이 아닌 "주장 문장"인가?
□ 평가 배점 높은 항목에 더 많은 페이지가 배분되었는가?
□ 수치·데이터가 있는 슬라이드는 number 또는 note 필드에 출처가 있는가?"""

    return call_json(prompt, max_tokens=8192)


# ─────────────────────────────────────────────
# PPTX 파일 생성
# ─────────────────────────────────────────────

def build_pptx(slide_data: dict, case: dict, progress_cb=None) -> bytes:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]   # Blank layout

    slides = slide_data.get("slides", [])
    total  = len(slides)
    meta   = {
        "client_name":  case.get("client_name", ""),
        "project_name": case.get("project_name", ""),
    }

    for i, sl in enumerate(slides):
        if progress_cb:
            progress_cb("PPTX 파일 생성 중...", i + 1, total)

        prs_slide = prs.slides.add_slide(blank)
        sl_type   = sl.get("type", "content")
        num       = sl.get("number", i + 1)

        if sl_type == "cover":
            _draw_cover(prs_slide, sl, num, total, meta)
        elif sl_type == "toc":
            _draw_toc(prs_slide, sl, num, total)
        elif sl_type == "process":
            _draw_process(prs_slide, sl, num, total)
        elif sl_type == "compare":
            _draw_compare(prs_slide, sl, num, total)
        elif sl_type == "number":
            _draw_number(prs_slide, sl, num, total)
        elif sl_type == "message":
            _draw_message(prs_slide, sl, num, total)
        else:
            # content (기본)
            _draw_content(prs_slide, sl, num, total)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _narrative_to_slide(ns: dict, idx: int) -> dict:
    """ppt_narrator 설계안 슬라이드 → ppt_generator 슬라이드 dict 변환."""
    sl_type = ns.get("slide_type", "content")
    num     = ns.get("number", idx + 1)
    head    = ns.get("head_copy", "")
    msg     = ns.get("key_message", "")
    ev      = ns.get("evidence", "")

    if sl_type == "cover":
        return {"number": num, "type": "cover", "title": head,
                "client": "", "date": ""}

    if sl_type == "toc":
        items = [line.lstrip("—-•· \t") for line in msg.splitlines() if line.strip()]
        return {"number": num, "type": "toc", "title": head, "items": items}

    if sl_type == "message":
        return {"number": num, "type": "message", "title": "",
                "message": head, "note": msg or ev}

    if sl_type == "number":
        lines = [l.strip() for l in msg.splitlines() if l.strip()]
        return {"number": num, "type": "number", "title": head,
                "value":  lines[0] if lines else "",
                "label":  lines[1] if len(lines) > 1 else "",
                "description": ev or (lines[2] if len(lines) > 2 else "")}

    if sl_type == "process":
        lines = [l.strip().lstrip("—-•· \t") for l in msg.splitlines() if l.strip()]
        steps = [{"label": l, "desc": ""} for l in lines[:5]]
        if not steps:
            steps = [{"label": head, "desc": ev}]
        return {"number": num, "type": "process", "title": head, "steps": steps}

    if sl_type == "compare":
        lines = [l.strip() for l in msg.splitlines() if l.strip()]
        mid   = len(lines) // 2
        left  = [l.lstrip("—-•· \t") for l in lines[:mid]] or ["(내용 없음)"]
        right = [l.lstrip("—-•· \t") for l in lines[mid:]] or ["(내용 없음)"]
        return {"number": num, "type": "compare", "title": head,
                "left_title":  "현재 문제",
                "left_items":  left,
                "right_title": "해결 방향 / INTERZ",
                "right_items": right}

    # content (기본)
    bullets = [l.strip().lstrip("—-•· \t") for l in msg.splitlines() if l.strip()]
    if not bullets and ev:
        bullets = [ev]
    return {"number": num, "type": "content", "title": head, "bullets": bullets or [msg]}


def build_pptx_from_narrative(narrative_slides: list, case: dict,
                               progress_cb=None) -> bytes:
    """ppt_narrator 설계안 → PPTX bytes (Claude 재호출 없음)."""
    import datetime
    today = datetime.date.today().strftime("%Y년 %m월")
    client = case.get("client_name", "")

    converted = []
    for i, ns in enumerate(narrative_slides):
        sl = _narrative_to_slide(ns, i)
        if sl["type"] == "cover":
            sl["client"] = client
            sl["date"]   = today
        converted.append(sl)

    slide_data = {"slides": converted}
    return build_pptx(slide_data, case, progress_cb)


# ─────────────────────────────────────────────
# 메인 진입점
# ─────────────────────────────────────────────

def run(case_detail: dict, pages: int, progress_cb=None) -> bytes:
    """PPTX bytes 반환."""
    slide_data = generate_slides(case_detail, pages, progress_cb)
    case       = case_detail.get("case", {})
    pptx_bytes = build_pptx(slide_data, case, progress_cb)

    if progress_cb:
        slides = slide_data.get("slides", [])
        progress_cb("완료!", len(slides), len(slides))

    return pptx_bytes
