# output/pptx_builder.py
# 역할: 최종 제안서 데이터를 PPTX 파일로 출력
# - 오케스트레이터의 final_proposal 데이터를 슬라이드로 변환
# - 슬라이드 구성: 표지 → 전략 → 컨셉 → 실행계획 → 대본 → 유통전략 → 회사소개 → Q&A
# - 색상/폰트/레이아웃 일관성 유지

import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt, Emu

from core.dna import ConceptDNA


# ─────────────────────────────────────────────
# 디자인 토큰 (모노크롬 — 흰 배경 / 검정 텍스트 / 컬러 없음)
# ─────────────────────────────────────────────

_C_NAVY   = RGBColor(0x00, 0x00, 0x00)   # 헤더 배경 → 검정
_C_GOLD   = RGBColor(0x33, 0x33, 0x33)   # 강조 → 진한 회색
_C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
_C_LIGHT  = RGBColor(0xF8, 0xF8, 0xF8)   # 본문 배경 → 거의 흰색
_C_TEXT   = RGBColor(0x00, 0x00, 0x00)   # 본문 텍스트 → 검정
_C_GRAY   = RGBColor(0x77, 0x77, 0x77)   # 보조 텍스트 → 회색
_C_ACCENT = RGBColor(0x44, 0x44, 0x44)   # 강조 포인트 → 어두운 회색

# 폰트 (Windows/Mac 공통 한글 지원)
_FONT_KO = "맑은 고딕"
_FONT_EN = "Calibri"

# 슬라이드 크기 (16:9 와이드)
_SLIDE_W = Cm(33.867)
_SLIDE_H = Cm(19.05)

# 공통 여백
_MARGIN_L = Cm(2.0)
_MARGIN_T = Cm(3.8)
_CONTENT_W = Cm(29.867)


# ─────────────────────────────────────────────
# RFP 기반 목차 (TOC)
# ─────────────────────────────────────────────

_DEFAULT_TOC = [
    "표지",
    "목차",
    "사업 이해 및 분석",
    "추진 전략",
    "핵심 컨셉",
    "콘텐츠 기획",
    "제작 계획",
    "대본/스토리보드",
    "유통/마케팅 전략",
    "추진 일정",
    "예산 계획",
    "수행 실적",
    "회사 소개",
    "마무리",
]


def _build_rfp_toc(rfp_data: dict) -> list:
    """RFP 분석 결과(rfp_analysis step)에서 제안서 목차 구성.

    core_tasks가 있으면 해당 내용을 기반으로 중간 섹션 구성.
    없으면 기본 목차(_DEFAULT_TOC) 반환.
    """
    if not rfp_data or not isinstance(rfp_data, dict):
        return _DEFAULT_TOC[:]

    core_tasks = rfp_data.get("core_tasks", [])
    if not core_tasks or not isinstance(core_tasks, list):
        return _DEFAULT_TOC[:]

    fixed_front = ["표지", "목차", "사업 이해 및 분석", "추진 전략", "핵심 컨셉"]
    fixed_back  = ["추진 일정", "예산 계획", "수행 실적", "회사 소개", "마무리"]
    existing    = set(fixed_front + fixed_back)

    task_sections = []
    for task in core_tasks:
        if isinstance(task, str):
            section = task.strip()
            if section and section not in existing:
                existing.add(section)
                task_sections.append(section)

    if not task_sections:
        return _DEFAULT_TOC[:]

    return fixed_front + task_sections + fixed_back


# ─────────────────────────────────────────────
# 진입점
# ─────────────────────────────────────────────

def build_pptx(dna: ConceptDNA, final_proposal: dict, output_dir: str = None,
               rfp_data: dict = None) -> Path:
    """제안서 전체를 PPTX로 생성.

    Args:
        dna: 최종 ConceptDNA (표지/메타 정보 사용)
        final_proposal: orchestrator의 final_proposal dict
        output_dir: 저장 경로 (None이면 output/proposals/ 사용)

    Returns:
        생성된 PPTX 파일 경로
    """
    prs = Presentation()
    prs.slide_width  = _SLIDE_W
    prs.slide_height = _SLIDE_H

    rfp_toc = _build_rfp_toc(rfp_data or {})

    # 섹션별 슬라이드 추가
    _add_cover_slide(prs, dna, final_proposal.get("cover", {}))
    _add_toc_slide(prs, dna, rfp_toc)
    _add_strategy_slide(prs, final_proposal.get("strategy", {}), dna)
    _add_concept_slide(prs, final_proposal.get("concept", {}), dna)
    _add_execution_slide(prs, final_proposal.get("episodes", {}), dna)
    _add_script_slide(prs, final_proposal.get("script", {}))
    _add_distribution_slide(prs, final_proposal.get("marketing", {}))
    _add_company_slide(prs, final_proposal.get("company", {}))
    _add_qa_slide(prs, final_proposal.get("qa", {}))

    # 파일 저장
    out_dir = Path(output_dir) if output_dir else Path(__file__).parent / "proposals"
    out_dir.mkdir(parents=True, exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = re.sub(r"[^\w가-힣\-]", "_", dna.project_name)[:40]
    filename = f"{safe_name}_{timestamp}.pptx"
    file_path = out_dir / filename

    prs.save(str(file_path))
    return file_path


# ─────────────────────────────────────────────
# 슬라이드별 구현
# ─────────────────────────────────────────────

def _add_cover_slide(prs: Presentation, dna: ConceptDNA, cover: dict) -> None:
    """표지 슬라이드: 발주처, 사업명, 컨셉, 슬로건."""
    slide = _blank_slide(prs)

    # 전체 배경 → 네이비
    _fill_background(slide, _C_NAVY)

    # 금색 상단 라인 (장식)
    _add_rect(slide, Cm(0), Cm(0), _SLIDE_W, Cm(0.5), _C_GOLD)

    # 발주처 (소제목)
    _add_text_box(
        slide, dna.client_name,
        Cm(2.5), Cm(2.5), Cm(28), Cm(1.2),
        font_size=18, bold=False, color=_C_GOLD, align=PP_ALIGN.LEFT,
    )

    # 사업명 (메인 타이틀)
    _add_text_box(
        slide, dna.project_name,
        Cm(2.5), Cm(4.0), Cm(28), Cm(3.5),
        font_size=32, bold=True, color=_C_WHITE, align=PP_ALIGN.LEFT,
        wrap=True,
    )

    # 슬로건
    slogan = dna.slogan or cover.get("slogan", "")
    if slogan:
        _add_text_box(
            slide, f'"{slogan}"',
            Cm(2.5), Cm(8.0), Cm(28), Cm(1.5),
            font_size=20, bold=False, color=_C_GOLD, align=PP_ALIGN.LEFT,
        )

    # 컨셉 한 줄
    concept = dna.concept or cover.get("concept", "")
    if concept:
        _add_text_box(
            slide, concept,
            Cm(2.5), Cm(9.8), Cm(28), Cm(1.2),
            font_size=14, bold=False, color=_C_GRAY, align=PP_ALIGN.LEFT,
        )

    # 하단 구분선
    _add_rect(slide, Cm(2.5), Cm(14.5), Cm(28.867), Cm(0.08), _C_GOLD)

    # 하단 정보 (영상 종류 / 수량 / 기간)
    meta_parts = [dna.video_type, f"{dna.quantity}편 / 편당 {dna.duration}"]
    if dna.deadline:
        meta_parts.append(f"납품기한: {dna.deadline}")
    _add_text_box(
        slide, "   |   ".join(p for p in meta_parts if p),
        Cm(2.5), Cm(14.8), Cm(24), Cm(0.9),
        font_size=11, bold=False, color=_C_GRAY, align=PP_ALIGN.LEFT,
    )

    # 제출 날짜
    today = datetime.now().strftime("%Y년 %m월")
    _add_text_box(
        slide, today,
        Cm(26), Cm(14.8), Cm(5.5), Cm(0.9),
        font_size=11, bold=False, color=_C_GRAY, align=PP_ALIGN.RIGHT,
    )


def _add_toc_slide(prs: Presentation, dna: ConceptDNA,
                   rfp_toc: list = None) -> None:
    """목차 슬라이드. rfp_toc가 있으면 RFP 기반 목차, 없으면 기본 목차 사용."""
    slide = _blank_slide(prs)
    _fill_background(slide, _C_LIGHT)
    _add_header_bar(slide, "목차  /  Contents")

    # 표지/목차/마무리 제외하고 표시 (최대 10개)
    toc = rfp_toc if rfp_toc else _DEFAULT_TOC
    display_items = [s for s in toc if s not in ("표지", "목차", "마무리")][:10]

    col_w = Cm(13.5)
    for idx, section in enumerate(display_items):
        num = f"{idx+1:02d}"
        row = idx // 2
        col = idx % 2
        x   = _MARGIN_L + col * (col_w + Cm(1.0))
        y   = Cm(4.2) + row * Cm(2.5)

        if y + Cm(2.0) > _SLIDE_H - Cm(1.0):
            break  # 슬라이드 하단 초과 방지

        # 번호 박스
        _add_rect(slide, x, y, Cm(1.4), Cm(1.4), _C_NAVY)
        _add_text_box(slide, num, x, y, Cm(1.4), Cm(1.4),
                      font_size=12, bold=True, color=_C_GOLD, align=PP_ALIGN.CENTER)

        # 타이틀
        _add_text_box(slide, section, x + Cm(1.6), y, col_w - Cm(1.8), Cm(1.4),
                      font_size=13, bold=True, color=_C_TEXT, align=PP_ALIGN.LEFT)


def _add_strategy_slide(prs: Presentation, strategy: dict, dna: ConceptDNA) -> None:
    """전략 슬라이드: 위기제시 → 현황진단 → 해결책."""
    # ── 슬라이드 1: 위기·현황 진단 ──
    slide = _blank_slide(prs)
    _fill_background(slide, _C_LIGHT)
    _add_header_bar(slide, "01  현황 분석 & 전략")

    blocks = [
        ("위기 제시", dna.crisis_statement or strategy.get("crisis_statement", ""), _C_ACCENT),
        ("현황 진단", dna.current_situation or strategy.get("current_situation", ""), _C_NAVY),
        ("해결책 방향", dna.solution_direction or strategy.get("solution_direction", ""), _C_GOLD),
    ]

    for i, (label, text, color) in enumerate(blocks):
        x = _MARGIN_L + i * Cm(9.8)
        y = Cm(4.0)
        # 색상 박스 헤더
        _add_rect(slide, x, y, Cm(9.2), Cm(0.9), color)
        _add_text_box(slide, label, x, y, Cm(9.2), Cm(0.9),
                      font_size=12, bold=True, color=_C_WHITE, align=PP_ALIGN.CENTER)
        # 내용 박스
        _add_rect(slide, x, y + Cm(0.9), Cm(9.2), Cm(9.5), _C_WHITE)
        _add_text_box(slide, text or "(내용 없음)",
                      x + Cm(0.3), y + Cm(1.2), Cm(8.6), Cm(8.8),
                      font_size=11, bold=False, color=_C_TEXT, align=PP_ALIGN.LEFT,
                      wrap=True)

    # ── 슬라이드 2: 평가항목 배점 분석 ──
    eval_items = dna.evaluation_items or strategy.get("evaluation_items", [])
    if eval_items:
        slide2 = _blank_slide(prs)
        _fill_background(slide2, _C_LIGHT)
        _add_header_bar(slide2, "01  평가항목 분석")

        headers = ["평가항목", "배점", "대응 전략 포인트"]
        rows = []
        for item in eval_items[:8]:
            if isinstance(item, dict):
                rows.append([
                    item.get("item", ""),
                    str(item.get("score", "")),
                    item.get("strategy", "") or item.get("description", ""),
                ])
            else:
                rows.append([str(item), "", ""])

        _add_table(slide2, headers, rows,
                   Cm(2.0), Cm(3.8), Cm(29.867), Cm(0.75))


def _add_concept_slide(prs: Presentation, concept: dict, dna: ConceptDNA) -> None:
    """컨셉 슬라이드: 핵심 컨셉, 슬로건, 톤앤매너."""
    slide = _blank_slide(prs)
    _fill_background(slide, _C_NAVY)
    _add_header_bar(slide, "02  핵심 컨셉 & 슬로건", header_bg=_C_GOLD, text_color=_C_NAVY)

    # 빅아이디어 박스
    big_idea = dna.concept or concept.get("concept", "")
    _add_text_box(
        slide, big_idea,
        Cm(2.5), Cm(4.0), Cm(28.0), Cm(2.5),
        font_size=28, bold=True, color=_C_WHITE, align=PP_ALIGN.CENTER, wrap=True,
    )

    # 슬로건 후보 3개
    slogans = dna.slogans or concept.get("slogans", [])
    if slogans:
        for i, s in enumerate(slogans[:3]):
            y = Cm(7.0) + i * Cm(2.2)
            text = s.get("text", s) if isinstance(s, dict) else str(s)
            rationale = s.get("rationale", "") if isinstance(s, dict) else ""
            star = "★" if i == 0 else "☆"
            _add_rect(slide, Cm(2.5), y, Cm(1.2), Cm(1.6), _C_GOLD if i == 0 else _C_ACCENT)
            _add_text_box(slide, star, Cm(2.5), y, Cm(1.2), Cm(1.6),
                          font_size=14, bold=True, color=_C_WHITE, align=PP_ALIGN.CENTER)
            _add_text_box(slide, text, Cm(4.0), y, Cm(25.5), Cm(0.85),
                          font_size=14 if i == 0 else 12, bold=(i == 0),
                          color=_C_GOLD if i == 0 else _C_GRAY,
                          align=PP_ALIGN.LEFT)
            if rationale:
                _add_text_box(slide, rationale, Cm(4.0), y + Cm(0.9), Cm(25.5), Cm(0.7),
                              font_size=10, bold=False, color=_C_GRAY, align=PP_ALIGN.LEFT)

    # 톤앤매너 + 감성 키워드
    slide2 = _blank_slide(prs)
    _fill_background(slide2, _C_LIGHT)
    _add_header_bar(slide2, "02  톤앤매너 & 비주얼 방향")

    tone = dna.tone_and_manner or concept.get("tone_and_manner", "")
    keywords = dna.tone_keywords or concept.get("tone_keywords", [])
    visual = dna.visual_direction or concept.get("visual_direction", "")

    _add_text_box(slide2, "[ 톤앤매너 ]",
                  Cm(2.0), Cm(3.8), Cm(10), Cm(0.8),
                  font_size=12, bold=True, color=_C_NAVY, align=PP_ALIGN.LEFT)
    _add_text_box(slide2, tone or "-",
                  Cm(2.0), Cm(4.7), Cm(14.5), Cm(4.0),
                  font_size=12, bold=False, color=_C_TEXT, align=PP_ALIGN.LEFT, wrap=True)

    # 감성 키워드 칩
    _add_text_box(slide2, "[ 감성 키워드 ]",
                  Cm(17.5), Cm(3.8), Cm(10), Cm(0.8),
                  font_size=12, bold=True, color=_C_NAVY, align=PP_ALIGN.LEFT)
    for i, kw in enumerate(keywords[:5]):
        kx = Cm(17.5) + (i % 3) * Cm(4.0)
        ky = Cm(4.7) + (i // 3) * Cm(1.5)
        _add_rect(slide2, kx, ky, Cm(3.6), Cm(1.1), _C_NAVY)
        _add_text_box(slide2, str(kw), kx, ky, Cm(3.6), Cm(1.1),
                      font_size=12, bold=False, color=_C_WHITE, align=PP_ALIGN.CENTER)

    # 비주얼 방향
    if visual:
        _add_text_box(slide2, "[ 비주얼 방향 ]",
                      Cm(2.0), Cm(9.2), Cm(10), Cm(0.8),
                      font_size=12, bold=True, color=_C_NAVY, align=PP_ALIGN.LEFT)
        _add_text_box(slide2, visual,
                      Cm(2.0), Cm(10.1), Cm(29.5), Cm(4.0),
                      font_size=11, bold=False, color=_C_TEXT, align=PP_ALIGN.LEFT, wrap=True)


def _add_execution_slide(prs: Presentation, episodes_data: dict, dna: ConceptDNA) -> None:
    """실행계획 슬라이드: 편별 제작 방향, 일정, 예산."""
    # ── 편별 제작 기획 ──
    slide = _blank_slide(prs)
    _fill_background(slide, _C_LIGHT)
    _add_header_bar(slide, "03  편별 실행 계획")

    episodes = dna.episodes or episodes_data.get("episodes", [])
    if episodes:
        headers = ["편수", "제목", "핵심 메시지", "주요 장면"]
        rows = []
        for ep in episodes[:6]:
            if isinstance(ep, dict):
                rows.append([
                    f"{ep.get('ep_num', '')}편",
                    ep.get("title", ""),
                    ep.get("key_message", ep.get("message", "")),
                    ep.get("key_scene", ep.get("scene", "")),
                ])
            else:
                rows.append([str(ep), "", "", ""])
        _add_table(slide, headers, rows,
                   Cm(2.0), Cm(3.8), Cm(29.867), Cm(0.75))

    # ── 제작 일정 ──
    schedule = dna.production_schedule or episodes_data.get("production_schedule", [])
    if schedule:
        slide2 = _blank_slide(prs)
        _fill_background(slide2, _C_LIGHT)
        _add_header_bar(slide2, "03  제작 일정")

        headers2 = ["단계", "기간", "주요 작업", "산출물"]
        rows2 = []
        for phase in schedule:
            if isinstance(phase, dict):
                rows2.append([
                    phase.get("phase", phase.get("stage", "")),
                    phase.get("period", phase.get("duration", "")),
                    phase.get("tasks", phase.get("work", "")),
                    phase.get("deliverable", phase.get("output", "")),
                ])
            else:
                rows2.append([str(phase), "", "", ""])
        _add_table(slide2, headers2, rows2,
                   Cm(2.0), Cm(3.8), Cm(29.867), Cm(0.75))

    # ── 예산 계획 ──
    budget_plan = dna.budget_plan or episodes_data.get("budget_plan", {})
    if budget_plan:
        slide3 = _blank_slide(prs)
        _fill_background(slide3, _C_LIGHT)
        _add_header_bar(slide3, "03  예산 배분 계획")

        items_raw = budget_plan.get("items", budget_plan.get("breakdown", []))
        if isinstance(items_raw, list) and items_raw:
            headers3 = ["항목", "금액", "비율", "비고"]
            rows3 = []
            for item in items_raw:
                if isinstance(item, dict):
                    rows3.append([
                        item.get("category", item.get("name", "")),
                        item.get("amount", ""),
                        item.get("ratio", item.get("percent", "")),
                        item.get("note", ""),
                    ])
                else:
                    rows3.append([str(item), "", "", ""])
            total = budget_plan.get("total", dna.budget or "")
            if total:
                rows3.append(["합계", str(total), "100%", ""])
            _add_table(slide3, headers3, rows3,
                       Cm(2.0), Cm(3.8), Cm(29.867), Cm(0.75))


def _add_script_slide(prs: Presentation, script_data: dict) -> None:
    """대본/장면구성 슬라이드."""
    slide = _blank_slide(prs)
    _fill_background(slide, _C_LIGHT)
    _add_header_bar(slide, "04  대본 & 장면 구성")

    scripts = script_data.get("scripts", [])
    if not scripts:
        _add_text_box(slide, "대본이 생성되지 않았습니다.",
                      _MARGIN_L, Cm(5.0), Cm(28), Cm(2.0),
                      font_size=13, bold=False, color=_C_GRAY, align=PP_ALIGN.CENTER)
        return

    for ep_script in scripts[:3]:   # 최대 3편까지
        ep_slide = _blank_slide(prs)
        _fill_background(ep_slide, _C_LIGHT)

        ep_num   = ep_script.get("ep_num", ep_script.get("episode", ""))
        ep_title = ep_script.get("title", "")
        header_label = f"04  대본 — {ep_num}편: {ep_title}" if ep_title else f"04  대본 — {ep_num}편"
        _add_header_bar(ep_slide, header_label)

        scenes = ep_script.get("scenes", [])
        if scenes:
            headers = ["씬", "장소/상황", "영상 구성", "나레이션/대사"]
            rows = []
            for sc in scenes[:12]:
                if isinstance(sc, dict):
                    rows.append([
                        str(sc.get("scene_no", sc.get("num", ""))),
                        sc.get("location", sc.get("place", sc.get("setting", ""))),
                        sc.get("video", sc.get("visual", "")),
                        sc.get("narration", sc.get("dialogue", sc.get("script", ""))),
                    ])
                else:
                    rows.append(["", "", "", str(sc)])
            _add_table(ep_slide, headers, rows,
                       Cm(1.5), Cm(3.8), Cm(30.867), Cm(0.65), font_size=9)
        else:
            full_text = ep_script.get("full_script", ep_script.get("content", ""))
            if full_text:
                _add_text_box(ep_slide, str(full_text)[:800],
                              _MARGIN_L, Cm(4.0), Cm(29.0), Cm(12.0),
                              font_size=10, bold=False, color=_C_TEXT,
                              align=PP_ALIGN.LEFT, wrap=True)


def _add_distribution_slide(prs: Presentation, marketing: dict) -> None:
    """유통/마케팅 슬라이드: 채널 전략, KPI."""
    # ── 채널 전략 ──
    slide = _blank_slide(prs)
    _fill_background(slide, _C_LIGHT)
    _add_header_bar(slide, "05  유통 & 마케팅 전략")

    strategy_text = marketing.get("distribution_strategy", "")
    channels = marketing.get("distribution_channels", [])

    if strategy_text:
        _add_text_box(slide, strategy_text,
                      _MARGIN_L, Cm(4.0), _CONTENT_W, Cm(2.5),
                      font_size=12, bold=False, color=_C_TEXT,
                      align=PP_ALIGN.LEFT, wrap=True)

    if channels:
        headers = ["채널", "주요 전략", "목표 KPI"]
        rows = []
        for ch in channels[:7]:
            if isinstance(ch, dict):
                rows.append([
                    ch.get("channel", ch.get("platform", "")),
                    ch.get("strategy", ch.get("content", "")),
                    ch.get("kpi", ch.get("target", "")),
                ])
            else:
                rows.append([str(ch), "", ""])
        y_offset = Cm(7.0) if strategy_text else Cm(4.0)
        _add_table(slide, headers, rows,
                   Cm(2.0), y_offset, Cm(29.867), Cm(0.7))

    # ── KPI 목표 ──
    kpi_targets = marketing.get("kpi_targets", [])
    if kpi_targets:
        slide2 = _blank_slide(prs)
        _fill_background(slide2, _C_LIGHT)
        _add_header_bar(slide2, "05  KPI 목표 & 성과 측정")

        headers2 = ["지표", "목표값", "측정 주기", "비고"]
        rows2 = []
        for kpi in kpi_targets[:8]:
            if isinstance(kpi, dict):
                rows2.append([
                    kpi.get("metric", kpi.get("name", "")),
                    kpi.get("target", kpi.get("value", "")),
                    kpi.get("period", kpi.get("cycle", "월별")),
                    kpi.get("note", ""),
                ])
            else:
                rows2.append([str(kpi), "", "", ""])
        _add_table(slide2, headers2, rows2,
                   Cm(2.0), Cm(3.8), Cm(29.867), Cm(0.75))

        reporting = marketing.get("reporting_system", "")
        if reporting:
            _add_text_box(slide2, f"▶ 성과 보고 체계:  {reporting}",
                          _MARGIN_L, Cm(14.0), _CONTENT_W, Cm(1.2),
                          font_size=11, bold=False, color=_C_ACCENT, align=PP_ALIGN.LEFT)


def _add_company_slide(prs: Presentation, company: dict) -> None:
    """회사 소개 슬라이드."""
    slide = _blank_slide(prs)
    _fill_background(slide, _C_LIGHT)
    _add_header_bar(slide, "06  회사 소개 & 포트폴리오")

    intro = company.get("intro", "")
    if intro:
        _add_text_box(slide, intro,
                      _MARGIN_L, Cm(4.0), _CONTENT_W, Cm(3.0),
                      font_size=12, bold=False, color=_C_TEXT,
                      align=PP_ALIGN.LEFT, wrap=True)

    # 실적 표
    achievements = company.get("achievements", [])
    if achievements:
        headers = ["연도", "사업명", "발주처", "영상 종류"]
        rows = []
        for a in achievements[:8]:
            if isinstance(a, dict):
                rows.append([
                    str(a.get("year", "")),
                    a.get("project", a.get("name", "")),
                    a.get("client", ""),
                    a.get("type", a.get("video_type", "")),
                ])
            else:
                rows.append(["", str(a), "", ""])
        y_off = Cm(7.2) if intro else Cm(4.0)
        _add_table(slide, headers, rows,
                   Cm(2.0), y_off, Cm(29.867), Cm(0.7))

    # 팀 구성
    team = company.get("team_composition", {})
    if team:
        slide2 = _blank_slide(prs)
        _fill_background(slide2, _C_LIGHT)
        _add_header_bar(slide2, "06  투입 인력 구성")

        headers2 = ["역할", "담당자", "주요 경력"]
        rows2 = []
        for role, info in team.items():
            if isinstance(info, dict):
                rows2.append([role, info.get("name", ""), info.get("career", "")])
            else:
                rows2.append([role, str(info), ""])
        _add_table(slide2, headers2, rows2,
                   Cm(2.0), Cm(3.8), Cm(29.867), Cm(0.8))


def _add_qa_slide(prs: Presentation, qa: dict) -> None:
    """Q&A 슬라이드: 심사위원 예상 질문 & 답변."""
    qa_list = qa.get("qa_list", qa.get("pairs", []))
    if not qa_list:
        return

    slide = _blank_slide(prs)
    _fill_background(slide, _C_NAVY)
    _add_header_bar(slide, "07  예상 Q&A", header_bg=_C_GOLD, text_color=_C_NAVY)

    for i, pair in enumerate(qa_list[:5]):
        y_base = Cm(3.8) + i * Cm(2.8)
        q_text = pair.get("question", pair.get("q", "")) if isinstance(pair, dict) else str(pair)
        a_text = pair.get("answer", pair.get("a", "")) if isinstance(pair, dict) else ""

        # Q 레이블
        _add_rect(slide, Cm(1.5), y_base, Cm(1.0), Cm(0.9), _C_GOLD)
        _add_text_box(slide, "Q", Cm(1.5), y_base, Cm(1.0), Cm(0.9),
                      font_size=12, bold=True, color=_C_NAVY, align=PP_ALIGN.CENTER)
        _add_text_box(slide, q_text, Cm(2.7), y_base, Cm(29.0), Cm(0.9),
                      font_size=11, bold=True, color=_C_WHITE, align=PP_ALIGN.LEFT)

        # A 레이블
        _add_rect(slide, Cm(1.5), y_base + Cm(1.0), Cm(1.0), Cm(0.9), _C_ACCENT)
        _add_text_box(slide, "A", Cm(1.5), y_base + Cm(1.0), Cm(1.0), Cm(0.9),
                      font_size=12, bold=True, color=_C_WHITE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, a_text, Cm(2.7), y_base + Cm(1.0), Cm(29.0), Cm(0.9),
                      font_size=11, bold=False, color=_C_WHITE,
                      align=PP_ALIGN.LEFT, wrap=True)


# ─────────────────────────────────────────────
# 공통 UI 헬퍼
# ─────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    """빈 레이아웃 슬라이드 추가."""
    blank_layout = prs.slide_layouts[6]   # index 6 = blank
    return prs.slides.add_slide(blank_layout)


def _fill_background(slide, color: RGBColor) -> None:
    """슬라이드 배경색 채우기."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(
    slide,
    text: str,
    header_bg: RGBColor = None,
    text_color: RGBColor = None,
) -> None:
    """슬라이드 상단 헤더 바 추가."""
    bg = header_bg or _C_NAVY
    tc = text_color or _C_WHITE
    _add_rect(slide, Cm(0), Cm(0), _SLIDE_W, Cm(3.2), bg)
    _add_text_box(slide, text,
                  Cm(2.0), Cm(0.8), Cm(28), Cm(1.8),
                  font_size=16, bold=True, color=tc, align=PP_ALIGN.LEFT)
    # 골드 하단 라인
    _add_rect(slide, Cm(0), Cm(3.2), _SLIDE_W, Cm(0.12), _C_GOLD)


def _add_rect(slide, x, y, w, h, color: RGBColor):
    """단색 직사각형 도형 추가."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()   # 테두리 없음
    return shape


def _add_text_box(
    slide,
    text: str,
    x, y, w, h,
    font_size: int = 12,
    bold: bool = False,
    color: RGBColor = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = False,
) -> None:
    """텍스트 박스 추가 (투명 배경)."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap

    p = tf.paragraphs[0]
    p.alignment = align

    run = p.add_run()
    run.text = str(text) if text else ""

    run.font.size     = Pt(font_size)
    run.font.bold     = bold
    run.font.color.rgb = color or _C_TEXT
    run.font.name     = _FONT_KO

    return txBox


def _add_table(
    slide,
    headers: list,
    rows: list,
    x, y, w,
    row_height,
    font_size: int = 10,
) -> None:
    """데이터 테이블 추가."""
    if not rows:
        return

    total_rows = len(rows) + 1   # 헤더 포함
    cols       = len(headers)

    table = slide.shapes.add_table(total_rows, cols, x, y, w, row_height * total_rows).table
    table.first_row = True

    # 컬럼 너비 균등 배분 (첫 열 좁게)
    col_widths = _calc_col_widths(headers, w, cols)
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = int(cw)

    # 헤더 행
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _C_NAVY
        _set_cell_text(cell, header, font_size=font_size, bold=True, color=_C_WHITE, align=PP_ALIGN.CENTER)

    # 데이터 행
    for ri, row in enumerate(rows):
        bg = _C_WHITE if ri % 2 == 0 else _C_LIGHT
        for ci, value in enumerate(row[:cols]):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            _set_cell_text(cell, str(value) if value else "",
                           font_size=font_size, bold=False, color=_C_TEXT,
                           align=PP_ALIGN.LEFT)


def _calc_col_widths(headers: list, total_w, n_cols: int) -> list:
    """헤더 이름 기준 컬럼 너비 비율 계산."""
    # 첫 컬럼이 번호/편수/단계 등 짧은 경우 좁게
    short_first = headers[0] in ("편수", "씬", "연도", "단계", "역할", "채널")
    if short_first and n_cols >= 3:
        ratios = [0.1] + [(0.9 / (n_cols - 1))] * (n_cols - 1)
    elif n_cols == 4:
        ratios = [0.15, 0.3, 0.3, 0.25]
    elif n_cols == 3:
        ratios = [0.25, 0.45, 0.30]
    else:
        ratios = [1.0 / n_cols] * n_cols

    return [total_w * r for r in ratios]


def _set_cell_text(
    cell,
    text: str,
    font_size: int = 10,
    bold: bool = False,
    color: RGBColor = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """테이블 셀 텍스트 설정."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text

    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color or _C_TEXT
    run.font.name      = _FONT_KO


# ─────────────────────────────────────────────
# 템플릿 기반 PPT 생성
# ─────────────────────────────────────────────

def _default_style() -> dict:
    """기본 디자인 토큰."""
    return {
        "bg":  (255, 255, 255),
        "hd":  (30,  30,  30),
        "ht":  (255, 255, 255),
        "bd":  (30,  30,  30),
        "ac":  (79,  70,  229),
        "tf":  "맑은 고딕",
        "bf":  "맑은 고딕",
        "ts":  24,
        "bs":  13,
    }


def _extract_style_from_pdf(pdf_bytes: bytes) -> dict:
    """PDF 파일에서 디자인 토큰 추출 (pdfplumber 사용)."""
    sty = _default_style()
    try:
        import pdfplumber
        import io as _io

        with pdfplumber.open(_io.BytesIO(pdf_bytes)) as pdf:
            for page in pdf.pages[:3]:
                chars = page.chars or []
                for char in chars[:1000]:
                    size = float(char.get("size") or 0)
                    fontname = char.get("fontname") or ""
                    color = char.get("non_stroking_color")

                    # 폰트명 정리 (PDF 내부 임베드 접두어 제거: ABCDEF+FontName)
                    clean_font = re.sub(r'^[A-Z]{6}\+', '', fontname).strip()
                    if clean_font:
                        if size >= 14:
                            sty["tf"] = clean_font
                            sty["ts"] = max(16, min(40, int(size)))
                        elif 8 <= size < 14:
                            sty["bf"] = clean_font
                            sty["bs"] = max(9, min(18, int(size)))

                    # 텍스트 색상 추출
                    if color is not None:
                        try:
                            if isinstance(color, (int, float)):
                                # 단일 값 = 그레이스케일 (0.0=검정, 1.0=흰색)
                                v = int((1.0 - float(color)) * 255)
                                c = (v, v, v)
                            elif isinstance(color, (list, tuple)) and len(color) == 3:
                                # RGB: 0~1 범위 또는 0~255 범위
                                def _norm(x):
                                    x = float(x)
                                    return int(x * 255) if x <= 1.0 else int(x)
                                c = (_norm(color[0]), _norm(color[1]), _norm(color[2]))
                            else:
                                c = None

                            if c:
                                r, g, b = c
                                is_white = r > 240 and g > 240 and b > 240
                                is_black = r < 30 and g < 30 and b < 30
                                if not is_white:
                                    if size >= 14 and not is_black:
                                        sty["hd"] = c
                                        sty["ac"] = c
                                        brightness = 0.299 * r + 0.587 * g + 0.114 * b
                                        sty["ht"] = (255, 255, 255) if brightness < 140 else (0, 0, 0)
                                    elif not is_black:
                                        sty["bd"] = c
                        except Exception:
                            pass
    except ImportError:
        print("  [PDF style] pdfplumber 미설치 — 기본 스타일 사용")
    except Exception as e:
        print(f"  [PDF style] 추출 실패: {e}")
    return sty


def _extract_style_from_hwp(hwp_bytes: bytes, is_hwpx: bool = False) -> dict:
    """HWP/HWPX 파일에서 디자인 토큰 추출."""
    sty = _default_style()
    if not is_hwpx:
        # 바이너리 HWP는 파싱 복잡 — 기본 스타일 반환
        print("  [HWP style] 바이너리 HWP는 스타일 추출 미지원 — 기본 스타일 사용")
        return sty
    try:
        import zipfile
        import io as _io
        import xml.etree.ElementTree as ET

        with zipfile.ZipFile(_io.BytesIO(hwp_bytes)) as zf:
            names = zf.namelist()
            # HWPX 구조: Contents/header.xml에 스타일 있음
            style_candidates = [n for n in names
                                 if any(k in n.lower() for k in ('header', 'styles', 'fontsset', 'charpr'))]
            for fname in style_candidates[:4]:
                try:
                    xml_bytes = zf.read(fname)
                    root = ET.fromstring(xml_bytes)
                    # 폰트 정보 탐색
                    for el in root.iter():
                        tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
                        if tag in ('fontRef', 'font', 'face'):
                            name = el.get('name') or el.get('val') or ''
                            if name and len(name) > 1:
                                sty['tf'] = name
                                sty['bf'] = name
                                break
                        # 색상 정보 탐색
                        if tag in ('charPr', 'textColor', 'color'):
                            color_str = el.get('color') or el.get('val') or ''
                            if color_str and len(color_str) >= 6:
                                try:
                                    hex_c = color_str.lstrip('#')[-6:]
                                    r = int(hex_c[0:2], 16)
                                    g = int(hex_c[2:4], 16)
                                    b = int(hex_c[4:6], 16)
                                    is_white = r > 240 and g > 240 and b > 240
                                    is_black = r < 30 and g < 30 and b < 30
                                    if not is_white and not is_black:
                                        sty['ac'] = (r, g, b)
                                        sty['hd'] = (r, g, b)
                                        brightness = 0.299 * r + 0.587 * g + 0.114 * b
                                        sty['ht'] = (255, 255, 255) if brightness < 140 else (0, 0, 0)
                                except ValueError:
                                    pass
                except Exception:
                    continue
    except Exception as e:
        print(f"  [HWPX style] 추출 실패: {e}")
    return sty


def _extract_template_style(tmpl_prs) -> dict:
    """PPTX 파일에서 디자인 토큰 + 레이아웃 패턴 추출.

    sty 키:
      bg, hd, ht, bd, ac, tf, bf, ts, bs  — 기본 토큰
      palette   — 발견된 컬러 목록 (최대 5개, 밝기순 정렬)
      ac2       — 보조 강조색 (팔레트 2번째)
      has_line  — 구분선 도형 존재 여부
      layout    — 추출된 레이아웃 타입 힌트 ('cover'|'section'|'two_col'|'plain')
    """
    sty = _default_style()
    sty["palette"]  = []
    sty["ac2"]      = sty["ac"]
    sty["has_line"] = False
    sty["layout"]   = "plain"

    def _safe_rgb(rgb_obj):
        try:
            return (rgb_obj[0], rgb_obj[1], rgb_obj[2])
        except Exception:
            return None

    def _brightness(c):
        return 0.299 * c[0] + 0.587 * c[1] + 0.114 * c[2]

    def _is_neutral(c):
        r, g, b = c
        is_white = r > 235 and g > 235 and b > 235
        is_black = r < 20 and g < 20 and b < 20
        is_gray  = abs(r - g) < 18 and abs(g - b) < 18
        return is_white or is_black or is_gray

    color_scores: dict = {}   # rgb_tuple → 면적(px²) 누적

    num_slides = len(tmpl_prs.slides)
    for slide in (tmpl_prs.slides or [])[:5]:
        # 배경색
        try:
            fill = slide.background.fill
            if fill.type is not None:
                c = _safe_rgb(fill.fore_color.rgb)
                if c:
                    sty["bg"] = c
        except Exception:
            pass

        shape_count = 0
        textbox_count = 0
        wide_shape_count = 0   # 가로로 긴 도형 (헤더 바 후보)

        for shape in slide.shapes:
            shape_count += 1
            # ── 도형 채움색 + 면적 누적 ──
            try:
                if hasattr(shape, "fill") and shape.fill.type == 1:
                    c = _safe_rgb(shape.fill.fore_color.rgb)
                    if c and not _is_neutral(c):
                        w = shape.width or 0
                        h = shape.height or 0
                        area = (w / 914400) * (h / 914400)   # cm²
                        color_scores[c] = color_scores.get(c, 0) + area

                        # 가로로 넓은 도형 → 헤더 바 후보
                        w_cm = w / 914400
                        h_cm = h / 914400
                        if w_cm > 20 and h_cm < 5:
                            wide_shape_count += 1
                            sty["hd"] = c
                            br = _brightness(c)
                            sty["ht"] = (255, 255, 255) if br < 140 else (0, 0, 0)
            except Exception:
                pass

            # ── 선/구분선 감지 ──
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    sty["has_line"] = True
                elif hasattr(shape, "width") and hasattr(shape, "height"):
                    w_cm = (shape.width or 0) / 914400
                    h_cm = (shape.height or 0) / 914400
                    if w_cm > 15 and h_cm < 0.3:
                        sty["has_line"] = True
            except Exception:
                pass

            # ── 텍스트 박스 카운트 + 폰트 추출 ──
            if shape.has_text_frame:
                textbox_count += 1
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            fn  = run.font.name
                            sz  = run.font.size.pt if run.font.size else None
                            clr = _safe_rgb(run.font.color.rgb)
                            if fn and sz:
                                if sz >= 18:
                                    sty["tf"] = fn
                                    sty["ts"] = max(int(sz), sty["ts"])
                                    if clr and not _is_neutral(clr):
                                        sty["bd"] = clr
                                elif sz >= 9:
                                    sty["bf"] = fn
                                    sty["bs"] = int(sz)
                                    if clr and not _is_neutral(clr):
                                        sty["bd"] = clr
                        except Exception:
                            pass

        # ── 레이아웃 타입 힌트 ──
        if wide_shape_count >= 1 and textbox_count >= 2:
            sty["layout"] = "section"
        elif textbox_count >= 4:
            sty["layout"] = "two_col"

    # ── 컬러 팔레트 정리 (면적 큰 순서) ──
    sorted_colors = sorted(color_scores.items(), key=lambda x: -x[1])
    palette = []
    for c, _ in sorted_colors:
        if len(palette) >= 5:
            break
        # 기존 팔레트와 너무 유사한 색 제외 (유클리드 거리 < 30)
        too_close = any(
            sum((c[i] - p[i]) ** 2 for i in range(3)) ** 0.5 < 30
            for p in palette
        )
        if not too_close:
            palette.append(c)
    sty["palette"] = palette

    if palette:
        # 가장 면적 큰 색 → 헤더 (이미 wide_shape으로 설정됐으면 유지)
        if sty["hd"] == _default_style()["hd"] and palette:
            sty["hd"] = palette[0]
            br = _brightness(palette[0])
            sty["ht"] = (255, 255, 255) if br < 140 else (0, 0, 0)
        # 2번째 색 → 강조색
        if len(palette) >= 2:
            sty["ac"]  = palette[1]
            sty["ac2"] = palette[0]
        else:
            sty["ac"]  = palette[0]
            sty["ac2"] = palette[0]

    print(f"  [Template] palette={sty['palette']} hd={sty['hd']} ac={sty['ac']} "
          f"layout={sty['layout']} has_line={sty['has_line']}")
    return sty


def _t_rgb(sty_tuple) -> RGBColor:
    return RGBColor(*sty_tuple)


def _template_bg(slide, color_tuple):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _t_rgb(color_tuple)


def _tmpl_box(slide, left_cm, top_cm, w_cm, h_cm, text, font, size_pt,
              bold=False, color=(0, 0, 0), align=PP_ALIGN.LEFT, wrap=True):
    tx = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(w_cm), Cm(h_cm))
    tf = tx.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _t_rgb(color)
    return tx


def _build_template_slides(detail: dict, target_pages: int = 20) -> list:
    """detail dict → 슬라이드 데이터 목록 (15~30장 목표).

    detail 구조: {"case": {...}, "steps": {"strategy": {...}, "creative": {...}, ...}}
    """
    case  = detail.get("case", {})
    steps = detail.get("steps", {})

    # steps가 없는 경우 detail 자체를 steps로 폴백 (레거시 호환)
    if not steps:
        steps = detail

    def _s(key): return steps.get(key) or {}
    def _sv(key, *fields):
        d = _s(key)
        for f in fields:
            v = d.get(f)
            if v:
                return str(v)
        return ""

    slides = []

    # ── 1. 표지 ──────────────────────────────────────────────
    slides.append({
        "type":  "cover",
        "title": case.get("project_name", "제안서"),
        "sub":   case.get("client_name", ""),
        "date":  datetime.now().strftime("%Y년 %m월"),
    })

    # ── 2. 목차 (RFP 기반) ───────────────────────────────────
    rfp_toc      = _build_rfp_toc(_s("rfp_analysis"))
    toc_sections = [s for s in rfp_toc if s not in ("표지", "목차")]

    slides.append({
        "type":  "toc",
        "title": "목차",
        "items": toc_sections,
    })

    # ── 3. 공모 개요 분석 ────────────────────────────────────
    rfp = _s("rfp_analysis")
    if rfp:
        items = []
        for label, field in [
            ("사업 목적", "purpose"), ("사업 개요", "summary"), ("예산 규모", "budget"),
            ("핵심 요구 사항", "key_requirements"), ("금지 사항", "forbidden_notes"),
        ]:
            v = rfp.get(field)
            if v:
                text = ", ".join(str(x) for x in v) if isinstance(v, list) else str(v)
                items.append((label, text[:300]))

        eval_items = rfp.get("evaluation_items") or []
        if eval_items:
            eval_text = "\n".join(
                f"• {e.get('item', str(e))} ({e.get('score', '')}점)"
                if isinstance(e, dict) else f"• {e}"
                for e in eval_items[:8]
            )
            items.append(("평가 배점표", eval_text[:400]))

        _append_section_chunks(slides, "공모 개요 분석", items, max_per_slide=3)

    # ── 4. 리서치 & 현황 분석 ────────────────────────────────
    res = _s("research")
    if res:
        items = []
        for label, field in [
            ("시장 현황", "market_overview"), ("경쟁 분석", "competitor_analysis"),
            ("최근 이슈", "recent_issues"), ("유사 사례", "similar_cases"),
            ("타깃 인사이트", "target_insight"), ("핵심 기회", "key_opportunity"),
        ]:
            v = res.get(field)
            if v:
                if isinstance(v, list):
                    text = "\n".join(f"• {str(x)[:150]}" for x in v[:5])
                else:
                    text = str(v)[:400]
                items.append((label, text))

        if items:
            _append_section_chunks(slides, "리서치 & 현황 분석", items, max_per_slide=2)

    # ── 5. 전략 방향 ─────────────────────────────────────────
    strat = _s("strategy")
    if strat:
        items = []
        for label, field in [
            ("핵심 문제 정의", "core_problem"),
            ("위기 제시",      "crisis_statement"),
            ("현황 진단",      "current_situation"),
            ("해결 방향",      "solution_direction"),
            ("전략 요약",      "strategy_summary"),
        ]:
            v = strat.get(field)
            if v:
                items.append((label, str(v)[:350]))

        _append_section_chunks(slides, "전략 방향 수립", items, max_per_slide=2)

        # 기대 효과 → process 슬라이드
        efx = strat.get("expected_effects") or []
        if isinstance(efx, list) and efx:
            slides.append({
                "type":  "process",
                "title": "기대 효과",
                "steps": [str(e)[:100] for e in efx[:5]],
            })

        # 설득 구조 → process 슬라이드
        persu = strat.get("persuasion_structure") or []
        if isinstance(persu, list) and persu:
            slides.append({
                "type":  "process",
                "title": "설득 구조",
                "steps": [str(p)[:100] for p in persu[:5]],
            })

    # ── 6. 핵심 컨셉 & 슬로건 ───────────────────────────────
    cre = _s("creative")
    if cre:
        # 컨셉 메인 슬라이드
        items = []
        for label, field in [
            ("핵심 컨셉", "concept"),
            ("컨셉 설명", "concept_description"),
            ("대표 슬로건", "confirmed_slogan"),
            ("톤앤매너", "tone_description"),
            ("비주얼 방향", "visual_direction"),
        ]:
            v = cre.get(field)
            if v:
                items.append((label, str(v)[:350]))
        _append_section_chunks(slides, "핵심 컨셉 & 슬로건", items, max_per_slide=2)

        # 슬로건 후보 슬라이드 (있을 경우)
        slogans = cre.get("slogans") or []
        if slogans:
            sl_items = []
            for i, s in enumerate(slogans[:6]):
                if isinstance(s, dict):
                    text = s.get("text", s.get("slogan", str(s)))
                    rationale = s.get("rationale", s.get("reason", ""))
                    sl_items.append((f"슬로건 {i+1}: {text}", str(rationale)[:200]))
                else:
                    sl_items.append((f"슬로건 {i+1}", str(s)[:200]))
            if sl_items:
                slides.append({"type": "section", "title": "슬로건 후보", "items": sl_items})

        # 감성 키워드
        kws = cre.get("tone_keywords") or []
        if kws:
            slides.append({
                "type": "keywords",
                "title": "감성 키워드 & 톤앤매너",
                "keywords": [str(k) for k in kws[:12]],
            })

    # ── 7. 편별 실행 기획 ────────────────────────────────────
    plan = _s("plan")
    if plan:
        eps = plan.get("episodes") or []
        if isinstance(eps, list) and eps:
            # 편 목록 개요 슬라이드
            ep_overview = []
            for ep in eps[:8]:
                if isinstance(ep, dict):
                    num   = ep.get("episode_number", ep.get("ep_num", ""))
                    title = ep.get("title", "")
                    msg   = ep.get("core_message", ep.get("key_message", ep.get("message", "")))
                    ep_overview.append((f"{num}편. {title}".strip(". "), str(msg)[:200]))
            if ep_overview:
                slides.append({"type": "section", "title": "편별 실행 기획 개요", "items": ep_overview})

            # 편당 상세 슬라이드 (최대 6편)
            for ep in eps[:6]:
                if not isinstance(ep, dict):
                    continue
                num   = ep.get("episode_number", ep.get("ep_num", ""))
                title = ep.get("title", "")
                items = []
                for label, field in [
                    ("핵심 메시지", "core_message"), ("key_message", "key_message"),
                    ("주요 장면", "key_scene"), ("제작 방향", "production_direction"),
                    ("타깃 반응", "target_response"),
                ]:
                    v = ep.get(field)
                    if v:
                        items.append((label, str(v)[:300]))
                        break  # 핵심 메시지는 하나만
                for label, field in [
                    ("주요 장면", "key_scene"), ("제작 방향", "production_direction"),
                    ("타깃 반응", "target_response"), ("예산", "budget"),
                ]:
                    v = ep.get(field)
                    if v:
                        items.append((label, str(v)[:300]))
                if items:
                    slides.append({
                        "type": "section",
                        "title": f"{num}편. {title}".strip(". "),
                        "items": items,
                    })

        # 제작 일정
        schedule = plan.get("production_schedule") or []
        if isinstance(schedule, list) and schedule:
            sch_items = []
            for phase in schedule[:8]:
                if isinstance(phase, dict):
                    ph = phase.get("phase", phase.get("stage", phase.get("step", "")))
                    pd = phase.get("period", phase.get("duration", ""))
                    task = phase.get("tasks", phase.get("work", phase.get("description", "")))
                    sch_items.append((f"{ph}  [{pd}]".strip("  []"), str(task)[:200]))
                else:
                    sch_items.append(("단계", str(phase)[:200]))
            if sch_items:
                slides.append({"type": "section", "title": "제작 일정", "items": sch_items})

        # 예산 계획
        budget = plan.get("budget_plan") or {}
        if isinstance(budget, dict) and budget:
            budget_items_raw = budget.get("items", budget.get("breakdown", []))
            if isinstance(budget_items_raw, list) and budget_items_raw:
                b_items = []
                for item in budget_items_raw[:8]:
                    if isinstance(item, dict):
                        cat = item.get("category", item.get("name", ""))
                        amt = item.get("amount", "")
                        note = item.get("note", "")
                        b_items.append((str(cat), f"{amt}  {note}".strip()))
                if b_items:
                    slides.append({"type": "section", "title": "예산 계획", "items": b_items})

    # ── 8. 대본 개요 ─────────────────────────────────────────
    scripts = steps.get("script") or []
    if isinstance(scripts, list) and scripts:
        for sc_row in scripts[:3]:
            ep_num = sc_row.get("episode_number", "")
            script_data = sc_row.get("script") or {}
            scenes = script_data.get("scenes", [])
            ep_title = script_data.get("title", sc_row.get("title", ""))

            if scenes:
                sc_items = []
                for sc in scenes[:6]:
                    if isinstance(sc, dict):
                        scene_no = sc.get("scene_no", sc.get("num", ""))
                        loc = sc.get("location", sc.get("place", sc.get("setting", "")))
                        narr = sc.get("narration", sc.get("dialogue", sc.get("script", "")))
                        sc_items.append((
                            f"씬 {scene_no}  {loc}".strip(),
                            str(narr)[:250],
                        ))
                if sc_items:
                    label = f"{ep_num}편 대본 구성"
                    if ep_title:
                        label += f" — {ep_title}"
                    slides.append({"type": "section", "title": label, "items": sc_items})
            else:
                # scenes 없으면 full_script 일부라도
                full = script_data.get("full_script", sc_row.get("full_text", ""))
                if full:
                    slides.append({
                        "type": "section",
                        "title": f"{ep_num}편 대본",
                        "items": [("대본", str(full)[:500])],
                    })

    # ── 9. 마케팅 & 유통 전략 ───────────────────────────────
    mkt = _s("marketing")
    if mkt:
        m_items = []
        for label, field in [
            ("유튜브 전략", "youtube_strategy"),
            ("SNS 전략", "sns_strategy"),
            ("인플루언서 전략", "influencer_strategy"),
            ("배포 플랫폼", "platforms"),
            ("KPI 목표", "kpi"),
            ("마케팅 예산", "marketing_budget"),
        ]:
            v = mkt.get(field)
            if v:
                if isinstance(v, list):
                    text = "\n".join(f"• {str(x)[:120]}" for x in v[:5])
                elif isinstance(v, dict):
                    text = "\n".join(f"• {k}: {str(val)[:80]}"
                                     for k, val in list(v.items())[:5])
                else:
                    text = str(v)[:400]
                m_items.append((label, text))

        _append_section_chunks(slides, "마케팅 & 유통 전략", m_items, max_per_slide=2)

        # KPI → diagram 슬라이드
        kpi_raw = mkt.get("kpi") or mkt.get("kpi_targets") or []
        if isinstance(kpi_raw, list) and kpi_raw:
            kpi_labels = []
            for k in kpi_raw[:6]:
                if isinstance(k, dict):
                    metric = k.get("metric", k.get("name", ""))
                    target = k.get("target", k.get("value", ""))
                    kpi_labels.append(f"{metric}\n{target}".strip())
                else:
                    kpi_labels.append(str(k)[:60])
            if kpi_labels:
                slides.append({
                    "type":   "diagram",
                    "title":  "KPI 목표",
                    "labels": kpi_labels,
                })

    # ── 10. 예상 Q&A ─────────────────────────────────────────
    final = _s("final_proposal")
    if final:
        qa_prep = final.get("qa_prep") or []
        if isinstance(qa_prep, list) and qa_prep:
            qa_items = []
            for qa in qa_prep[:8]:
                if isinstance(qa, dict):
                    q = qa.get("question", qa.get("q", ""))
                    a = qa.get("answer", qa.get("a", ""))
                    if q:
                        qa_items.append((f"Q. {q[:120]}", str(a)[:300]))
            if qa_items:
                _append_section_chunks(slides, "예상 Q&A", qa_items, max_per_slide=3)

        # 회사 소개
        company = final.get("company_profile") or {}
        if isinstance(company, dict) and company:
            c_items = []
            for label, field in [
                ("회사 소개", "intro"),
                ("주요 실적", "achievements"),
                ("핵심 인력", "key_personnel"),
            ]:
                v = company.get(field)
                if v:
                    if isinstance(v, list):
                        text = "\n".join(f"• {str(x)[:120]}" for x in v[:5])
                    else:
                        text = str(v)[:350]
                    c_items.append((label, text))
            if c_items:
                slides.append({"type": "section", "title": "회사 소개", "items": c_items})

    # ── 11. 종료 ─────────────────────────────────────────────
    slides.append({
        "type":  "end",
        "title": "감사합니다",
        "sub":   case.get("client_name", ""),
    })

    print(f"  [Template] 슬라이드 {len(slides)}장 구성 (target={target_pages})")
    return slides


def _append_section_chunks(slides: list, title: str, items: list,
                             max_per_slide: int = 3):
    """items를 max_per_slide 단위로 나눠 여러 슬라이드로 분할."""
    if not items:
        return
    for i in range(0, len(items), max_per_slide):
        chunk = items[i:i + max_per_slide]
        suffix = f" ({i // max_per_slide + 1})" if len(items) > max_per_slide else ""
        slides.append({
            "type":  "section",
            "title": title + suffix,
            "items": chunk,
        })


def generate_from_template(detail: dict, template_bytes: bytes,
                            pages: int = 20, progress_cb=None,
                            file_ext: str = ".pptx") -> bytes:
    """참고 파일의 디자인 스타일을 적용한 제안서 PPT 생성.

    Args:
        detail: get_case_detail() 반환값
        template_bytes: 업로드된 파일 바이트 (.pptx / .pdf / .hwp / .hwpx)
        pages: 최대 슬라이드 수 (현재는 content 기반으로 자동 결정)
        progress_cb: callback(message, current, total)
        file_ext: 파일 확장자 (소문자, 점 포함 — 추출 방법 결정에 사용)

    Returns:
        생성된 PPTX bytes
    """
    import io as _io

    def _prog(msg, cur, tot):
        if progress_cb:
            try:
                progress_cb(msg, cur, tot)
            except Exception:
                pass

    ext = file_ext.lower()
    _prog("템플릿 스타일 분석 중...", 1, 5)

    if ext == ".pptx":
        tmpl_prs = Presentation(_io.BytesIO(template_bytes))
        sty = _extract_template_style(tmpl_prs)
    elif ext == ".pdf":
        sty = _extract_style_from_pdf(template_bytes)
    elif ext == ".hwpx":
        sty = _extract_style_from_hwp(template_bytes, is_hwpx=True)
    else:  # .hwp (바이너리) — 기본 스타일
        sty = _extract_style_from_hwp(template_bytes, is_hwpx=False)

    print(f"  [Template/{ext}] bg={sty['bg']} hd={sty['hd']} ac={sty['ac']} "
          f"font={sty['tf']}/{sty['bf']}")

    _prog("슬라이드 내용 구성 중...", 2, 5)
    slides_data = _build_template_slides(detail, target_pages=pages)

    _prog("PPT 파일 생성 중...", 3, 5)
    prs = Presentation()
    prs.slide_width  = _SLIDE_W
    prs.slide_height = _SLIDE_H
    blank = prs.slide_layouts[6]

    SW = _SLIDE_W.cm
    SH = _SLIDE_H.cm

    def _hdr_bar(sl, title):
        """상단 헤더 바 + 제목 텍스트."""
        h = sl.shapes.add_shape(1, Cm(0), Cm(0), Cm(SW), Cm(2.4))
        h.fill.solid(); h.fill.fore_color.rgb = _t_rgb(sty["hd"])
        h.line.fill.background()
        _tmpl_box(sl, 1.5, 0.4, SW - 3, 1.7, title,
                  sty["tf"], max(sty["ts"], 14), bold=True, color=sty["ht"])

    for idx, sdata in enumerate(slides_data):
        sl = prs.slides.add_slide(blank)
        st = sdata["type"]
        _template_bg(sl, sty["bg"])

        if st == "cover":
            # 상단 2/5 색상 블록
            hdr = sl.shapes.add_shape(1, Cm(0), Cm(0), Cm(SW), Cm(SH * 0.55))
            hdr.fill.solid(); hdr.fill.fore_color.rgb = _t_rgb(sty["hd"])
            hdr.line.fill.background()
            _tmpl_box(sl, 2.0, 1.2, SW - 4, SH * 0.35,
                      sdata["title"], sty["tf"],
                      min(sty["ts"] + 8, 38), bold=True, color=sty["ht"], wrap=True)
            _tmpl_box(sl, 2.0, SH * 0.55 + 0.6, SW - 4, 1.6,
                      sdata["sub"], sty["bf"], sty["bs"] + 3, color=sty["bd"])
            _tmpl_box(sl, 2.0, SH * 0.55 + 2.4, SW - 4, 1.2,
                      sdata["date"], sty["bf"], sty["bs"], color=sty["bd"])

        elif st == "toc":
            _hdr_bar(sl, sdata["title"])
            items = sdata.get("items", [])
            n = max(len(items), 1)
            row_h = (SH - 3.0) / n
            y = 2.8
            for i, sec in enumerate(items):
                num_box = sl.shapes.add_shape(1, Cm(1.5), Cm(y), Cm(1.0), Cm(row_h * 0.7))
                num_box.fill.solid(); num_box.fill.fore_color.rgb = _t_rgb(sty["ac"])
                num_box.line.fill.background()
                _tmpl_box(sl, 1.5, y, Cm(1.0).cm, row_h * 0.7, str(i + 1).zfill(2),
                          sty["bf"], sty["bs"], bold=True, color=sty["ht"], align=PP_ALIGN.CENTER)
                _tmpl_box(sl, 3.2, y + 0.05, SW - 5, row_h * 0.7,
                          sec, sty["bf"], sty["bs"] + 1, color=sty["bd"])
                y += row_h

        elif st == "end":
            _template_bg(sl, sty["hd"])
            _tmpl_box(sl, 0, SH / 2 - 3, SW, 4.5,
                      sdata["title"], sty["tf"], 40,
                      bold=True, color=sty["ht"], align=PP_ALIGN.CENTER)
            if sdata.get("sub"):
                _tmpl_box(sl, 0, SH / 2 + 1.8, SW, 2,
                          sdata["sub"], sty["bf"], sty["bs"] + 3,
                          color=sty["ht"], align=PP_ALIGN.CENTER)

        elif st == "section":
            _hdr_bar(sl, sdata["title"])
            items = sdata.get("items", [])
            n = max(len(items), 1)
            avail_h = SH - 3.0
            row_h = avail_h / n
            y = 2.8
            for label, content in items:
                label_h = 0.85
                content_h = max(row_h - label_h - 0.2, 1.0)
                # 라벨 (강조색 배경)
                lb = sl.shapes.add_shape(1, Cm(1.5), Cm(y), Cm(SW - 3), Cm(label_h))
                lb.fill.solid(); lb.fill.fore_color.rgb = _t_rgb(sty["ac"])
                lb.line.fill.background()
                _tmpl_box(sl, 1.6, y + 0.07, SW - 3.2, label_h - 0.14,
                          label, sty["bf"], sty["bs"],
                          bold=True, color=sty["ht"])
                # 내용
                _tmpl_box(sl, 1.5, y + label_h + 0.1, SW - 3, content_h,
                          content, sty["bf"], sty["bs"], color=sty["bd"], wrap=True)
                y += row_h

        elif st == "bullets":
            _hdr_bar(sl, sdata["title"])
            items = sdata.get("items", [])
            n = max(len(items), 1)
            row_h = (SH - 3.0) / n
            y = 2.8
            for item in items:
                _tmpl_box(sl, 1.5, y, SW - 3, row_h,
                          f"• {item}", sty["bf"], sty["bs"],
                          color=sty["bd"], wrap=True)
                y += row_h

        elif st == "keywords":
            _hdr_bar(sl, sdata["title"])
            keywords = sdata.get("keywords", [])
            cols = 4
            kw_w = (SW - 3.0) / cols
            for i, kw in enumerate(keywords[:12]):
                col = i % cols
                row = i // cols
                kx = 1.5 + col * kw_w
                ky = 3.0 + row * 2.2
                box = sl.shapes.add_shape(1, Cm(kx), Cm(ky), Cm(kw_w - 0.3), Cm(1.7))
                box.fill.solid(); box.fill.fore_color.rgb = _t_rgb(sty["hd"])
                box.line.fill.background()
                _tmpl_box(sl, kx, ky, kw_w - 0.3, 1.7,
                          str(kw), sty["bf"], sty["bs"] + 1,
                          color=sty["ht"], align=PP_ALIGN.CENTER)

        elif st == "process":
            # 단계별 흐름도: 번호 배지 + 화살표 심볼 + 내용
            _hdr_bar(sl, sdata["title"])
            steps_list = sdata.get("steps", [])
            n = max(len(steps_list), 1)
            # 가로 배치 (최대 5단계) vs 세로 배치
            if n <= 5:
                item_w = (SW - 3.0) / n
                y_top  = 3.2
                item_h = SH - y_top - 0.5
                for i, step_text in enumerate(steps_list):
                    x = 1.5 + i * item_w
                    # 번호 배지
                    badge = sl.shapes.add_shape(1, Cm(x + item_w / 2 - 0.6), Cm(y_top),
                                                Cm(1.2), Cm(1.2))
                    badge.fill.solid()
                    badge.fill.fore_color.rgb = _t_rgb(sty["ac"])
                    badge.line.fill.background()
                    _tmpl_box(sl, x + item_w / 2 - 0.6, y_top, 1.2, 1.2,
                              str(i + 1), sty["tf"], sty["bs"] + 2,
                              bold=True, color=sty["ht"], align=PP_ALIGN.CENTER)
                    # 내용 박스
                    cont_box = sl.shapes.add_shape(1, Cm(x + 0.1), Cm(y_top + 1.5),
                                                   Cm(item_w - 0.3), Cm(item_h - 1.7))
                    cont_box.fill.solid()
                    cont_box.fill.fore_color.rgb = _t_rgb(
                        sty["ac2"] if i % 2 == 0 else sty["hd"])
                    cont_box.line.fill.background()
                    _tmpl_box(sl, x + 0.2, y_top + 1.6, item_w - 0.5, item_h - 1.9,
                              step_text, sty["bf"], sty["bs"],
                              color=sty["ht"], wrap=True, align=PP_ALIGN.CENTER)
                    # 화살표 (마지막 제외)
                    if i < n - 1:
                        _tmpl_box(sl, x + item_w - 0.5, y_top + item_h / 2 - 0.4,
                                  0.9, 0.9, "▶", sty["tf"], sty["bs"] + 4,
                                  color=sty["ac"], align=PP_ALIGN.CENTER)
            else:
                # 세로 배치 (5개 초과)
                row_h = (SH - 3.2) / n
                y = 3.2
                for i, step_text in enumerate(steps_list):
                    badge = sl.shapes.add_shape(1, Cm(1.5), Cm(y + 0.1),
                                                Cm(1.0), Cm(row_h - 0.2))
                    badge.fill.solid()
                    badge.fill.fore_color.rgb = _t_rgb(sty["ac"])
                    badge.line.fill.background()
                    _tmpl_box(sl, 1.5, y + 0.1, 1.0, row_h - 0.2,
                              str(i + 1), sty["bf"], sty["bs"],
                              bold=True, color=sty["ht"], align=PP_ALIGN.CENTER)
                    _tmpl_box(sl, 3.0, y + 0.1, SW - 4.5, row_h - 0.2,
                              step_text, sty["bf"], sty["bs"],
                              color=sty["bd"], wrap=True)
                    y += row_h

        elif st == "diagram":
            # 카드 그리드 다이어그램 (KPI 등)
            _hdr_bar(sl, sdata["title"])
            labels = sdata.get("labels", [])
            n = len(labels)
            if n == 0:
                pass
            else:
                cols = min(n, 3)
                rows = (n + cols - 1) // cols
                card_w = (SW - 3.0) / cols
                card_h = min((SH - 3.5) / rows, 4.0)
                for i, lbl in enumerate(labels):
                    col = i % cols
                    row = i // cols
                    cx = 1.5 + col * card_w
                    cy = 3.2 + row * (card_h + 0.3)
                    card = sl.shapes.add_shape(1, Cm(cx), Cm(cy),
                                               Cm(card_w - 0.3), Cm(card_h))
                    card.fill.solid()
                    card.fill.fore_color.rgb = _t_rgb(sty["hd"])
                    card.line.fill.background()
                    # 상단 강조 띠
                    accent_bar = sl.shapes.add_shape(1, Cm(cx), Cm(cy),
                                                     Cm(card_w - 0.3), Cm(0.4))
                    accent_bar.fill.solid()
                    accent_bar.fill.fore_color.rgb = _t_rgb(sty["ac"])
                    accent_bar.line.fill.background()
                    _tmpl_box(sl, cx + 0.2, cy + 0.5, card_w - 0.7, card_h - 0.6,
                              lbl, sty["bf"], sty["bs"] + 1,
                              color=sty["ht"], wrap=True, align=PP_ALIGN.CENTER)

            # 구분선 (참고 파일에 선이 있었으면)
            if sty.get("has_line"):
                line = sl.shapes.add_shape(1, Cm(1.5), Cm(2.6), Cm(SW - 3), Cm(0.06))
                line.fill.solid(); line.fill.fore_color.rgb = _t_rgb(sty["ac"])
                line.line.fill.background()

    _prog("파일 저장 중...", 4, 5)
    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    _prog("완료!", 5, 5)
    return buf.getvalue()


# ─────────────────────────────────────────────
# Gamma API 연동
# ─────────────────────────────────────────────

def has_gamma_key() -> bool:
    """GAMMA_API_KEY 환경변수 설정 여부 반환."""
    import os
    return bool(os.environ.get("GAMMA_API_KEY", "").strip())


def generate_with_gamma(topic: str, pages: int) -> dict:
    """Gamma API로 프레젠테이션 생성.

    공식 문서: https://public-api.gamma.app/v1.0/generations
    인증: X-API-KEY 헤더 (Bearer 아님)
    응답: {"generationId": "..."} → polling으로 완료 대기

    Args:
        topic: 제안서 내용 요약 (Gamma에 전달할 주제 텍스트)
        pages: 목표 슬라이드 수

    Returns:
        {
            "url":      str,        # Gamma 프레젠테이션 웹 URL
            "pptx_url": str | None, # 익스포트 URL (있을 경우)
        }

    Raises:
        RuntimeError: API 키 미설정 또는 API 오류 시
    """
    import os
    import time
    import requests

    api_key = os.environ.get("GAMMA_API_KEY", "").strip()

    # ── 진단 로그 ─────────────────────────────────
    print(f"  [Gamma] GAMMA_API_KEY: {'SET (len=%d)' % len(api_key) if api_key else 'NOT SET'}")

    if not api_key:
        raise RuntimeError(
            "GAMMA_API_KEY 환경변수가 설정되지 않았습니다.\n"
            "Gamma 가입 후 Settings → API Keys에서 키를 발급받아 "
            "Railway Variables에 GAMMA_API_KEY=발급받은키 로 추가하세요."
        )

    num_cards = max(5, min(50, pages))

    # ── 공식 헤더: X-API-KEY (Bearer 아님) ──────────
    headers = {
        "X-API-KEY":    api_key,
        "Content-Type": "application/json",
    }

    # ── Step 1: 생성 요청 ─────────────────────────
    # 올바른 base URL: public-api.gamma.app (api.gamma.app 아님)
    _BASE = "https://public-api.gamma.app"
    print(f"  [Gamma] POST {_BASE}/v1.0/generations — numCards={num_cards}")
    resp = requests.post(
        f"{_BASE}/v1.0/generations",
        json={
            "inputText": topic[:8000],
            "textMode":  "generate",
            "format":    "presentation",
            "numCards":  num_cards,
            "exportAs":  "pdf",      # pdf: 공식 문서 지원 확인된 포맷
        },
        headers=headers,
        timeout=60,
    )

    print(f"  [Gamma] HTTP {resp.status_code} — {resp.text[:300]}")

    if resp.status_code == 401:
        raise RuntimeError("Gamma API 인증 실패 — GAMMA_API_KEY를 확인하세요.")
    if resp.status_code == 402:
        raise RuntimeError("Gamma 플랜 문제 — gamma.app에서 플랜을 확인하세요.")
    if not resp.ok:
        raise RuntimeError(
            f"Gamma API 오류 ({resp.status_code}): {resp.text[:400]}"
        )

    data = resp.json()
    # 초기 응답은 generationId만 반환: {"generationId": "..."}
    print(f"  [Gamma] 초기 응답 keys={list(data.keys())}")

    gen_id = data.get("generationId") or data.get("id")
    if not gen_id:
        raise RuntimeError(
            f"Gamma API가 generationId를 반환하지 않았습니다. 응답: {str(data)[:400]}"
        )

    # ── Step 2: 폴링 — 항상 실행 (초기 응답에 status 없음) ──
    print(f"  [Gamma] 생성 대기 중 (generationId={gen_id})...")
    for attempt in range(60):           # 최대 5분 (5s × 60)
        time.sleep(5)
        poll = requests.get(
            f"{_BASE}/v1.0/generations/{gen_id}",
            headers=headers,
            timeout=30,
        )
        if not poll.ok:
            print(f"  [Gamma] 폴링 오류 ({poll.status_code}: {poll.text[:100]}) — 재시도")
            continue
        poll_data = poll.json()
        poll_status = poll_data.get("status", "")
        print(f"  [Gamma] 폴링 {attempt+1}/60 — status={poll_status!r}, keys={list(poll_data.keys())}")
        if poll_status == "completed":
            data = poll_data
            break
        if poll_status in ("failed", "error", "cancelled"):
            raise RuntimeError(
                f"Gamma 생성 실패 (status={poll_status}): "
                f"{poll_data.get('error') or str(poll_data)[:200]}"
            )
    else:
        raise RuntimeError("Gamma 생성 시간 초과 (5분) — 나중에 다시 시도하세요.")

    # ── Step 3: URL 추출 ──────────────────────────
    presentation_url = (
        data.get("gammaUrl") or
        data.get("url") or
        data.get("view_url") or
        ""
    )
    export_url = (
        data.get("exportUrl") or
        data.get("export_url") or
        data.get("download_url") or
        None
    )

    print(f"  [Gamma] 완료 — gammaUrl={presentation_url[:80] if presentation_url else 'None'}")
    print(f"  [Gamma] exportUrl={export_url[:80] if export_url else 'None'}")

    if not presentation_url and not export_url:
        raise RuntimeError(
            f"Gamma API 응답에 URL이 없습니다. 전체 응답: {str(data)[:400]}"
        )

    return {"url": presentation_url, "pptx_url": export_url}
